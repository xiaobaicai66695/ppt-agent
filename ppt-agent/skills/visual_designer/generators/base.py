"""PPT generator base module - shared utilities and palette definitions."""
from __future__ import annotations

import os
from copy import deepcopy
from typing import Literal

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Palette definitions (浅色基调，饱和度 5%-45%，亮度 80%-95%)
# ---------------------------------------------------------------------------

PALETTES: dict[str, dict[str, str]] = {
    "ocean_soft": {
        "name": "雾霾蓝",
        "primary": "5A8AA8",
        "secondary": "7BA3B8",
        "accent": "A8C4D4",
        "background": "F5F8FA",
        "text": "3A4A52",
        "text_muted": "7A8A92",
        "light_bg": "E8F0F5",
        "divider": "C8D8E5",
    },
    "sage_calm": {
        "name": "鼠尾草绿",
        "primary": "6A8E85",
        "secondary": "84B59F",
        "accent": "A8C5BD",
        "background": "FAF5F2",
        "text": "3A4A46",
        "text_muted": "8A9A92",
        "light_bg": "E8F0EC",
        "divider": "C8D8D0",
    },
    "warm_terracotta": {
        "name": "陶土橙",
        "primary": "C47060",
        "secondary": "D4A574",
        "accent": "E8DDD0",
        "background": "FAF5F2",
        "text": "4A3A32",
        "text_muted": "9A8A82",
        "light_bg": "F0E8E0",
        "divider": "D8C8B8",
    },
    "charcoal_light": {
        "name": "浅炭灰",
        "primary": "5A6A75",
        "secondary": "8A9AA5",
        "accent": "C8D0D5",
        "background": "F8F8F8",
        "text": "3A4248",
        "text_muted": "8A9298",
        "light_bg": "E8ECEE",
        "divider": "C0C8CC",
    },
    "berry_cream": {
        "name": "玫瑰灰粉",
        "primary": "8D5A6B",
        "secondary": "B89098",
        "accent": "E0D0D4",
        "background": "FAF5F5",
        "text": "4A3238",
        "text_muted": "9A8288",
        "light_bg": "F0E4E8",
        "divider": "D8C0C8",
    },
    "lavender_mist": {
        "name": "薰衣草灰",
        "primary": "7A6A8A",
        "secondary": "9A8AAA",
        "accent": "C8B8D4",
        "background": "F8F5FA",
        "text": "3A3248",
        "text_muted": "8A7A98",
        "light_bg": "EDE4F4",
        "divider": "D0C0E0",
    },
    "forest_moss": {
        "name": "苔藓绿",
        "primary": "5A8E5A",
        "secondary": "7EBF7E",
        "accent": "A8D4A8",
        "background": "F5FAF5",
        "text": "1A2A1A",
        "text_muted": "5A7A5A",
        "light_bg": "E4F4E4",
        "divider": "A8C8A8",
    },
    "sunset_peach": {
        "name": "杏色",
        "primary": "C4A080",
        "secondary": "D8B898",
        "accent": "E8D8C8",
        "background": "FAF5F2",
        "text": "3A2010",
        "text_muted": "8A6A4A",
        "light_bg": "F0E4D8",
        "divider": "D8C8B0",
    },

    # ── 中国场景配色 ──
    "government_red": {
        "name": "政务红",
        "primary": "8B1A1A",
        "secondary": "B87070",
        "accent": "D4A0A0",
        "background": "FDF8F5",
        "text": "3A1A1A",
        "text_muted": "8A6A6A",
        "light_bg": "F5E8E5",
        "divider": "D8C0B8",
    },
    "patriotic_blue": {
        "name": "爱国蓝",
        "primary": "2B5797",
        "secondary": "4A7AB0",
        "accent": "A0B8D4",
        "background": "F5F8FC",
        "text": "1A2A48",
        "text_muted": "5A6A8A",
        "light_bg": "E4EAF4",
        "divider": "B8C8D8",
    },
    "civic_gold": {
        "name": "公民金",
        "primary": "B8860B",
        "secondary": "C8A040",
        "accent": "E8D890",
        "background": "FAF8F0",
        "text": "3A2A0A",
        "text_muted": "8A7A4A",
        "light_bg": "F0E8D8",
        "divider": "D8C8A0",
    },
    "debate_purple": {
        "name": "答辩紫",
        "primary": "5A4A7A",
        "secondary": "8A74AA",
        "accent": "C8B8D4",
        "background": "F8F5FA",
        "text": "2A1A3A",
        "text_muted": "7A6A8A",
        "light_bg": "EDE4F4",
        "divider": "C8B8D8",
    },
    "activity_orange": {
        "name": "活力橙",
        "primary": "C06030",
        "secondary": "C89060",
        "accent": "E8C8A0",
        "background": "FAF5F2",
        "text": "3A2010",
        "text_muted": "8A6A4A",
        "light_bg": "F0E0D0",
        "divider": "D8C0A8",
    },
    "report_green": {
        "name": "报告绿",
        "primary": "3A6A4A",
        "secondary": "5A8A6A",
        "accent": "A0C0B0",
        "background": "F5F8F5",
        "text": "1A2A1A",
        "text_muted": "5A7A5A",
        "light_bg": "E4F0E8",
        "divider": "A8C0B0",
    },
    "simple_gray": {
        "name": "简约灰",
        "primary": "4A4A4A",
        "secondary": "6A6A6A",
        "accent": "C0C0C0",
        "background": "F8F8F8",
        "text": "2A2A2A",
        "text_muted": "8A8A8A",
        "light_bg": "E8E8E8",
        "divider": "C0C0C0",
    },
}


def rgb(hex_color: str) -> RGBColor:
    """Convert hex string like '5A8AA8' to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# ---------------------------------------------------------------------------
# Presentation helpers
# ---------------------------------------------------------------------------

def new_presentation(
    width_in: float = 13.333,
    height_in: float = 7.5,
    palette: str = "ocean_soft",
) -> Presentation:
    """Create a new 16:9 presentation. No slide is created — generators add slides."""
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def set_slide_background(slide, palette: str = "ocean_soft"):
    """Set solid background color for a slide."""
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(colors["background"])


def set_image_background(slide, image_path: str, brightness: float = 0.95, palette: str = "ocean_soft"):
    """
    为幻灯片设置图片背景。

    Args:
        slide: python-pptx Slide 对象
        image_path: 背景图片的完整路径（支持 JPG/PNG）
        brightness: 亮度调整 (0.0-1.0)，值越小背景越暗。推荐值 0.8-0.95

    Example:
        set_image_background(slide, "D:/path/to/background.jpg")
    """
    if not os.path.exists(image_path):
        print(f"Warning: Background image not found: {image_path}")
        return

    img = Image.open(image_path)
    if img.mode != 'RGB':
        img = img.convert('RGB')

    if brightness < 1.0:
        import numpy as np
        arr = np.array(img).astype(np.float32)
        arr = np.clip(arr * brightness, 0, 255).astype(np.uint8)
        img = Image.fromarray(arr)

    # 调整为幻灯片尺寸 16:9 (像素为单位)
    target_w, target_h = 1920, 1080
    img_ratio = img.width / img.height
    target_ratio = target_w / target_h

    if img_ratio > target_ratio:
        new_h = target_h
        new_w = int(new_h * img_ratio)
    else:
        new_w = target_w
        new_h = int(new_w / img_ratio)

    img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
    left = (new_w - target_w) // 2
    top = (new_h - target_h) // 2
    img = img.crop((left, top, left + target_w, top + target_h))

    # 临时保存处理后的图片
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        temp_path = tmp.name
    img.save(temp_path, quality=95)

    # 添加图片到 slide 的 relationships，获取 rId
    img_shape = slide.shapes.add_picture(temp_path, 0, 0)
    r_id = None
    for rid, rel in slide.part._rels.items():
        if 'image' in rel.reltype.lower():
            r_id = rid
            break

    # 构建原生 <p:bg> blipFill XML（背景图片在所有形状之下）
    from lxml import etree
    bgPr = etree.Element(qn("p:bgPr"))
    blipFill = etree.SubElement(bgPr, qn("a:blipFill"))
    blipFill.set("dpi", "0")
    blip = etree.SubElement(blipFill, qn("a:blip"))
    blip.set(qn("r:embed"), r_id)
    stretch = etree.SubElement(blipFill, qn("a:stretch"))
    etree.SubElement(stretch, qn("a:fillRect"))
    etree.SubElement(bgPr, qn("a:effectLst"))

    bg = etree.Element(qn("p:bg"))
    bg.append(bgPr)

    # 替换 p:cSld 中的 background
    p_cSld = slide._element.find(qn("p:cSld"))
    existing_bg = p_cSld.find(qn("p:bg"))
    if existing_bg is not None:
        p_cSld.remove(existing_bg)
    p_cSld.insert(0, bg)

    # 移除之前添加的 picture shape（背景现在由 p:bg 承载）
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if hasattr(shape, 'shape_type') and shape.shape_type == 13:
            sp_tree.remove(shape.element)

    os.remove(temp_path)

    # 自动添加磨砂玻璃叠加层，提升内容可读性
    # add_frosted_glass_overlay 返回适合玻璃背景的颜色映射
    return add_frosted_glass_overlay(slide, palette)


def add_frosted_glass_overlay(slide, palette: str = "ocean_soft"):
    """
    在幻灯片上添加磨砂玻璃叠加层并返回适合的颜色映射。

    磨砂玻璃覆盖整张幻灯片，顶层所有文字和色块印在上面。
    返回的颜色映射中 accent/divider 等替换为适合玻璃背景的白色/深色，
    避免和背景图片的蓝色调违和。

    Returns:
        dict: 颜色映射，key 为 palette 颜色名，value 为适合玻璃背景的新颜色。
              例如 {"primary": "FFFFFF", "divider": "D0D0D0", "accent": "FFFFFF"}
    """
    add_rect(
        slide,
        left=0.0, top=0.0, width=13.333, height=7.5,
        fill_color=(255, 255, 255, 128),
        line_color=None,
    )
    # 玻璃模式下，文字和色块用深色系（在白色磨砂底上可读），
    # 填充色用深蓝色系（保持海洋配色风格，在磨砂玻璃上有足够对比度）
    glass_palette = PALETTES.get(palette, PALETTES["ocean_soft"]).copy()
    glass_palette.update({
        # 文字色 → 深色（在白色磨砂底上可读）
        "primary": "2D3748",   # 深灰（标题/小标题）
        "secondary": "4A5568",  # 中灰（小标题/kicker）
        "accent": "2C5282",     # 深蓝（强调）
        "text": "1A202C",      # 近黑（正文）
        "text_muted": "718096", # 中灰（次要文字）
        # 填充色 → 深蓝（保持海洋风格又有对比度）
        "divider": "A0B4C0",
        "light_bg": "C8D8E5",
        "background": "D8E8F0",
        # 象限色（SWOT等用原始色块时，深蓝比浅蓝对比度更好）
        "primary_fill": "3D6B82",
        "secondary_fill": "5A8AA8",
        "accent_fill": "7BA3B8",
    })
    return glass_palette


def save_presentation(prs: Presentation, output_path: str):
    """Save the presentation to a file path."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)


def _clone_xml(el):
    """Clone an lxml element by round-tripping through string serialization.
    This produces a clean copy detached from the source document tree,
    avoiding namespace conflicts that deepcopy can cause across different Presentations.
    """
    from lxml import etree

    return etree.fromstring(etree.tostring(el, encoding="unicode"))


def save_slide(slide, output_path: str):
    """
    Save a single slide as a standalone PPTX file.
    Creates a new presentation, clones the slide (including chart parts and
    background image resources) into it, and saves.
    """
    from copy import deepcopy
    from lxml import etree
    from pptx.oxml.ns import qn
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    src_slide_part = slide.part
    src_slide_rel = src_slide_part._rels

    new_prs = Presentation()
    new_prs.slide_width = src_slide_part.package.presentation_part.presentation.slide_width
    new_prs.slide_height = src_slide_part.package.presentation_part.presentation.slide_height

    blank_layout = new_prs.slide_layouts[6]
    dest_slide = new_prs.slides.add_slide(blank_layout)
    dest_slide_part = dest_slide.part

    # Remove placeholder shapes from the blank slide
    for shape in list(dest_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # ------------------------------------------------------------------
    # Copy background: handle both solid-color and image (blipFill) backgrounds.
    # ------------------------------------------------------------------
    p_cSld = dest_slide._element.find(qn("p:cSld"))
    # Remove any existing background
    existing_bg = p_cSld.find(qn("p:bg"))
    if existing_bg is not None:
        p_cSld.remove(existing_bg)

    src_p_cSld = slide._element.find(qn("p:cSld"))
    src_bg = src_p_cSld.find(qn("p:bg")) if src_p_cSld is not None else None

    if src_bg is not None:
        bg_clone = deepcopy(src_bg)
        # If it's a blipFill (image) background, copy the image part
        bgPr = bg_clone.find(qn("p:bgPr"))
        if bgPr is not None:
            blipFill = bgPr.find(qn("a:blipFill"))
            if blipFill is not None:
                blip = blipFill.find(qn("a:blip"))
                if blip is not None:
                    old_rid = blip.get(qn("r:embed"))
                    if old_rid and old_rid in src_slide_rel:
                        try:
                            img_part = src_slide_rel[old_rid].target_part
                            new_rid = dest_slide_part.relate_to(img_part, RT.IMAGE)
                            blip.set(qn("r:embed"), new_rid)
                        except Exception:
                            pass
        p_cSld.insert(0, bg_clone)

    # ------------------------------------------------------------------
    # Copy chart relationships.
    # ------------------------------------------------------------------
    for rel_id, rel in src_slide_rel.items():
        if rel.reltype == RT.CHART:
            chart_part = rel.target_part
            dest_slide_part.relate_to(chart_part, RT.CHART)

    # ------------------------------------------------------------------
    # Clone all shape elements.
    # Update chart rIds to point to the newly-added chart parts.
    # ------------------------------------------------------------------
    dest_spTree = dest_slide.shapes._spTree
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for shape in slide.shapes:
        el = deepcopy(shape.element)
        chart_el = el.find(qn("c:chart"))
        if chart_el is not None:
            old_rid = chart_el.get("{" + r_ns + "}embed")
            if old_rid and old_rid in src_slide_rel:
                src_rel = src_slide_rel[old_rid]
                if src_rel.reltype == RT.CHART:
                    for new_rel_id, new_rel in dest_slide_part.rels.items():
                        if new_rel.target_part == src_rel.target_part:
                            chart_el.set("{" + r_ns + "}embed", new_rel_id)
                            break
        dest_spTree.append(el)

    save_presentation(new_prs, output_path)


def add_source_line(
    slide,
    source: str,
    palette: str = "ocean_soft",
):
    """在幻灯片底部渲染数据来源/参考信息。
    当 source 非空时，在底部添加灰色小字来源行。
    """
    if not source:
        return

    # 底部淡灰色分割线
    add_rect(
        slide,
        left=0.5, top=6.85, width=12.0, height=0.01,
        fill_color="text_muted", palette=palette,
    )

    # 来源文字（小号灰色）
    add_text(
        slide,
        text=source,
        left=0.5, top=6.9, width=12.0, height=0.35,
        font_size=9, bold=False,
        color="text_muted", alignment="left",
        palette=palette,
    )


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    alignment: Literal["left", "center", "right"] = "left",
    palette: str = "ocean_soft",
    font_name: str = None,
    italic: bool = False,
    colors: dict = None,
) -> "pptx.shapes.shapetree.Shape":
    """Add a text box to the slide with consistent styling."""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic

    if colors is not None:
        p.font.color.rgb = rgb(colors.get(color, color))
    else:
        palette_colors = PALETTES.get(palette, PALETTES["ocean_soft"])
        p.font.color.rgb = rgb(palette_colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)

    if font_name:
        p.font.name = font_name

    return txbox


def add_paragraph(
    tf,
    text: str,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    palette: str = "ocean_soft",
    alignment: Literal["left", "center", "right"] = "left",
    space_before: float = 0,
    space_after: float = 0,
    italic: bool = False,
    font_name: str = None,
    colors: dict = None,
):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if colors is not None:
        p.font.color.rgb = rgb(colors.get(color, color))
    else:
        palette_colors = PALETTES.get(palette, PALETTES["ocean_soft"])
        p.font.color.rgb = rgb(palette_colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)
    if font_name:
        p.font.name = font_name
    return p


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add a filled rectangle.

    Args:
        fill_color: either a palette key (str) like "primary", or an (R, G, B, A) tuple.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    if isinstance(fill_color, tuple) and len(fill_color) == 4:
        r, g, b, a = fill_color
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        # Set alpha transparency directly in XML: 0=opaque, 100000=fully transparent
        from pptx.oxml.ns import qn
        spPr = shape.element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(int(a / 255 * 100000)))
    else:
        fill.solid()
        fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_round_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    fill.solid()
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_ellipse(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add an ellipse/circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    fill.solid()
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "primary",
    width: float = 1.5,
    palette: str = "ocean_soft",
) -> "pptx.shapes.shapetree.Shape":
    """Add a straight line connector."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = rgb(PALETTES.get(palette, PALETTES["ocean_soft"]).get(color, color))
    connector.line.width = Pt(width)
    return connector


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "secondary",
    width: float = 1.5,
    palette: str = "ocean_soft",
) -> "pptx.shapes.shapetree.Shape":
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1), Inches(y1), Inches(x2 - x1), Inches(y2 - y1),
    )
    shape.fill.solid()
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    shape.fill.fore_color.rgb = rgb(colors.get(color, color))
    shape.line.fill.background()
    return shape


def add_text_in_shape(
    shape,
    text: str,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    palette: str = "ocean_soft",
    alignment: Literal["left", "center", "right"] = "left",
    vertical_alignment: str = "top",
    font_name: str = None,
    italic: bool = False,
):
    """Set text inside an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = True

    # Clear existing paragraphs except first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    p.font.color.rgb = rgb(colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)

    if font_name:
        p.font.name = font_name

    # Vertical alignment
    tf.anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
        "center": MSO_ANCHOR.MIDDLE,
    }.get(vertical_alignment, MSO_ANCHOR.TOP)


# ---------------------------------------------------------------------------
# Background utilities
# ---------------------------------------------------------------------------

def get_background_path(
    theme: str = None,
    scenario: str = None,
    random_select: bool = False,
) -> str:
    """
    获取背景图片路径。

    Args:
        theme: 主题标识 (如 "party_government")
        scenario: 场景关键词 (如 "党建汇报")
        random_select: 是否随机选择

    Returns:
        背景图片路径，或 None
    """
    from .background_manager import get_background as _get_bg
    return _get_bg(theme=theme, scenario=scenario, random_select=random_select)


def resolve_background(background: str) -> str:
    """
    解析 background 参数，返回实际图片路径。

    支持三种格式：
    - 绝对路径: "D:/path/to/image.jpg"
    - 主题名: "party_government"
    - 场景描述: "党建汇报"

    Args:
        background: background 参数值

    Returns:
        实际图片路径，或 None
    """
    if not background:
        return None

    # 已经是绝对路径
    if os.path.isabs(background) and os.path.exists(background):
        return background

    # 文件存在
    if os.path.exists(background):
        return background

    # 尝试作为 theme 或 scenario
    # theme 匹配时：固定使用该主题，随机选其中一张图片
    # scenario 匹配时：取优先级最高的匹配
    theme_bg = get_background_path(theme=background)
    if theme_bg:
        return get_background_path(theme=background, random_select=True)
    return get_background_path(scenario=background)
