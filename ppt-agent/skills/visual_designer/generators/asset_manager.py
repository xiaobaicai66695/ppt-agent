"""Offline visual asset helpers for visual_designer generators."""

from __future__ import annotations

from functools import lru_cache
import json
from pathlib import Path
import re
from typing import Iterable

from PIL import Image
from pptx.util import Inches

from .base import add_rect, set_image_background


ASSET_ROOT = Path(__file__).resolve().parents[1] / "assets"
MANIFEST_PATH = ASSET_ROOT / "manifest.json"
SUPPORTED_TYPES = {"icon", "background", "photo", "pattern"}
EXTERNAL_METADATA_FIELDS = (
    "source_id",
    "source_url",
    "download_url",
    "license",
    "attribution",
    "dimensions",
)


@lru_cache(maxsize=1)
def load_manifest_data() -> dict:
    if not MANIFEST_PATH.exists():
        return {"version": 0, "assets": []}
    data = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    return data if isinstance(data, dict) else {"version": 0, "assets": []}


def load_manifest() -> list[dict]:
    assets = load_manifest_data().get("assets", [])
    return assets if isinstance(assets, list) else []


def asset_record(asset_id: str) -> dict | None:
    for asset in load_manifest():
        if asset.get("id") == asset_id:
            return asset
    return None


def asset_path(asset_id: str) -> str | None:
    asset = asset_record(asset_id)
    if not asset:
        return None
    path = ASSET_ROOT / str(asset.get("path", ""))
    return str(path) if path.is_file() else None


def validate_manifest() -> list[str]:
    """Return portable manifest errors without depending on process cwd."""
    errors: list[str] = []
    data = load_manifest_data()
    if data.get("version") != 2:
        errors.append(f"unsupported_version:{data.get('version', '<missing>')}")

    source_ids = {
        str(source.get("id", "")).strip()
        for source in data.get("sources", [])
        if isinstance(source, dict)
    }
    seen: set[str] = set()
    asset_root = ASSET_ROOT.resolve()
    for asset in load_manifest():
        asset_id = str(asset.get("id", "")).strip()
        relative = str(asset.get("path", "")).strip()
        kind = str(asset.get("type", "")).strip()
        if not asset_id or asset_id in seen:
            errors.append(f"duplicate_or_empty_id:{asset_id or '<empty>'}")
            continue
        seen.add(asset_id)
        if kind not in SUPPORTED_TYPES:
            errors.append(f"unsupported_type:{asset_id}:{kind or '<empty>'}")
        if not relative:
            errors.append(f"empty_path:{asset_id}")
            continue

        target = (ASSET_ROOT / relative).resolve()
        try:
            target.relative_to(asset_root)
        except ValueError:
            errors.append(f"path_outside_bundle:{asset_id}:{relative}")
            continue
        if not target.is_file():
            errors.append(f"missing:{asset_id}:{relative}")
            continue
        _validate_case(asset_id, relative, errors)

        source_id = str(asset.get("source_id", "")).strip()
        if not source_id:
            errors.append(f"missing_source_id:{asset_id}")
        elif source_id not in source_ids and not source_id.startswith("project-"):
            errors.append(f"unknown_source_id:{asset_id}:{source_id}")
        if source_id and not source_id.startswith("project-"):
            for field in EXTERNAL_METADATA_FIELDS:
                if not asset.get(field):
                    errors.append(f"missing_external_metadata:{asset_id}:{field}")

        expected_dimensions = asset.get("dimensions")
        if not (
            isinstance(expected_dimensions, list)
            and len(expected_dimensions) == 2
            and all(isinstance(value, int) and value > 0 for value in expected_dimensions)
        ):
            errors.append(f"invalid_dimensions:{asset_id}")
            continue
        try:
            with Image.open(target) as image:
                actual_dimensions = [image.width, image.height]
                if actual_dimensions != expected_dimensions:
                    errors.append(
                        f"dimension_mismatch:{asset_id}:"
                        f"{expected_dimensions[0]}x{expected_dimensions[1]}:"
                        f"{image.width}x{image.height}"
                    )
                if kind == "icon" and "A" not in image.getbands():
                    errors.append(f"icon_without_alpha:{asset_id}")
        except Exception as exc:
            errors.append(f"unreadable_image:{asset_id}:{type(exc).__name__}")
    return errors


def _validate_case(asset_id: str, relative: str, errors: list[str]) -> None:
    cursor = ASSET_ROOT
    for part in Path(relative).parts:
        try:
            names = {child.name for child in cursor.iterdir()}
        except OSError:
            return
        if part not in names:
            errors.append(f"case_mismatch:{asset_id}:{relative}")
            return
        cursor = cursor / part


def find_asset(kind: str, tags: Iterable[str] = (), preferred_id: str = "") -> dict | None:
    assets = [asset for asset in load_manifest() if asset.get("type") == kind]
    if preferred_id:
        preferred = next((asset for asset in assets if asset.get("id") == preferred_id), None)
        if preferred:
            return preferred
    tag_set = {_normalize_text(tag) for tag in tags if tag}
    if tag_set:
        ranked: list[tuple[int, int, str, dict]] = []
        for asset in assets:
            candidates = {
                _normalize_text(str(value))
                for value in (
                    *asset.get("tags", []),
                    *asset.get("keywords", []),
                    *asset.get("recommended_templates", []),
                )
                if value
            }
            score = len(tag_set.intersection(candidates))
            if score:
                ranked.append((score, int(asset.get("priority", 0)), str(asset.get("id", "")), asset))
        if ranked:
            ranked.sort(reverse=True)
            return ranked[0][3]
    return assets[0] if assets else None


def resolve_asset_background(background: str | None, role: str = "") -> str | None:
    if background and background.startswith("asset:"):
        return asset_path(background.split(":", 1)[1])
    if background:
        return asset_path(background)
    asset = find_asset("background", tags=[role], preferred_id=default_background_for_role(role))
    return asset_path(str(asset.get("id", ""))) if asset else None


def default_background_for_role(role: str) -> str:
    return {
        "title": "editorial_workspace",
        "section": "editorial_city",
        "hero": "editorial_creative",
        "quote": "editorial_heritage",
        "summary": "editorial_nature",
    }.get(role, "")


def apply_asset_background(
    slide,
    background: str | None,
    palette: str,
    role: str,
    brightness: float = 0.96,
):
    path = resolve_asset_background(background, role)
    if path:
        return set_image_background(slide, path, brightness=brightness, palette=palette)
    return None


def _asset_id_from_text(text: str, kind: str, fallback: str = "") -> str:
    """Resolve one asset kind from manifest keywords with deterministic ranking."""
    normalized = _normalize_text(text)
    direct_id = normalized.replace("-", "_").replace(" ", "_")
    direct = asset_record(direct_id)
    if direct and direct.get("type") == kind:
        return direct_id

    ranked: list[tuple[int, int, int, int, str]] = []
    for index, asset in enumerate(load_manifest()):
        if asset.get("type") != kind:
            continue
        asset_id = str(asset.get("id", ""))
        candidates = [asset_id, *asset.get("keywords", []), *asset.get("tags", [])]
        best_length = 0
        match_count = 0
        for candidate in candidates:
            keyword = _normalize_text(str(candidate))
            if keyword and _keyword_matches(normalized, keyword):
                best_length = max(best_length, len(keyword))
                match_count += 1
        if best_length:
            ranked.append((best_length, match_count, int(asset.get("priority", 0)), -index, asset_id))
    if ranked:
        ranked.sort(reverse=True)
        return ranked[0][4]
    fallback_asset = asset_record(fallback)
    return fallback if fallback_asset and fallback_asset.get("type") == kind else ""


def icon_id_from_text(text: str, fallback: str = "") -> str:
    """Resolve the best icon from manifest keywords; return empty when unknown."""
    return _asset_id_from_text(text, "icon", fallback)


def photo_id_from_text(text: str, fallback: str = "photo_business_work") -> str:
    """Resolve a content photo from slide semantics with a stable general fallback."""
    return _asset_id_from_text(text, "photo", fallback)


def resolve_photo(image_path: str = "", text: str = "", fallback: str = "photo_business_work") -> str | None:
    """Resolve an explicit local image or a registered semantic content photo."""
    explicit = (image_path or "").strip()
    if explicit:
        asset_id = explicit.split(":", 1)[1] if explicit.startswith("asset:") else explicit
        record = asset_record(asset_id)
        if record and record.get("type") in {"photo", "background"}:
            registered = asset_path(asset_id)
            if registered:
                return registered
        candidate = Path(explicit).expanduser()
        if candidate.is_file():
            return str(candidate.resolve())

    photo_id = photo_id_from_text(text, fallback=fallback)
    return asset_path(photo_id) if photo_id else None


def add_cropped_photo(
    slide,
    image_path: str,
    left: float,
    top: float,
    width: float,
    height: float,
):
    """Add a center-cropped, replaceable PowerPoint picture to a fixed frame."""
    path = Path(image_path)
    if not path.is_file() or width <= 0 or height <= 0:
        return None

    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = width / height

    picture = slide.shapes.add_picture(
        str(path),
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    if image_ratio > frame_ratio:
        visible = frame_ratio / image_ratio
        picture.crop_left = picture.crop_right = max(0.0, (1.0 - visible) / 2.0)
    elif image_ratio < frame_ratio:
        visible = image_ratio / frame_ratio
        picture.crop_top = picture.crop_bottom = max(0.0, (1.0 - visible) / 2.0)
    picture.name = f"Replaceable photo - {path.stem}"
    return picture


def _normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", (value or "").strip().lower())


def _keyword_matches(text: str, keyword: str) -> bool:
    if re.search(r"[\u3400-\u9fff]", keyword):
        return keyword in text
    return bool(re.search(rf"(?<![a-z0-9]){re.escape(keyword)}(?![a-z0-9])", text))


def add_local_icon(
    slide,
    icon_id: str,
    left: float,
    top: float,
    size: float,
    palette: str = "ocean_soft",
    with_badge: bool = False,
    badge_color: str = "light_bg",
):
    """Add a registered icon, or omit it when the id is unknown."""
    path = asset_path(icon_id)
    if not path:
        return None
    if with_badge:
        add_rect(
            slide,
            left=left - size * 0.08,
            top=top - size * 0.08,
            width=size * 1.16,
            height=size * 1.16,
            fill_color=badge_color,
            palette=palette,
            line_color="divider",
            line_width=0.5,
        )
    return slide.shapes.add_picture(
        path,
        Inches(left),
        Inches(top),
        width=Inches(size),
        height=Inches(size),
    )


def add_pattern_overlay(
    slide,
    pattern_id: str,
    left: float,
    top: float,
    width: float,
    height: float,
    opacity_backdrop: bool = True,
    palette: str = "ocean_soft",
):
    path = asset_path(pattern_id)
    if not path:
        return None
    if opacity_backdrop:
        add_rect(
            slide,
            left=left,
            top=top,
            width=width,
            height=height,
            fill_color=(255, 255, 255, 170),
            palette=palette,
        )
    return slide.shapes.add_picture(
        path,
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
