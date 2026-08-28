import unittest
from io import BytesIO
from pathlib import Path
import re
import sys
import tempfile
import zipfile
from unittest.mock import patch

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT))

from generators import render_task
from generators import new_presentation, save_presentation, save_slide, set_image_background
from generators.base import PALETTES, background_image_palette, compact_source_text
from generators.comparison_table_generator import generate as generate_comparison_table
from generators.kpi_dashboard_generator import generate as generate_kpi_dashboard
from generators.component_layout import (
    choose_card_columns,
    clamp_text,
    component_accent,
    component_body,
    component_items,
    component_label,
    image_text_regions,
    render_component_slide,
)


class RenderTaskComponentsTest(unittest.TestCase):
    def test_extract_cards_prefers_semantic_components(self):
        plan = {
            "summary": "三层能力矩阵支撑端到端交付",
            "components": [
                {"id": "headline_1", "type": "headline", "text": "三层能力矩阵支撑端到端交付"},
                {"id": "feature_card_1", "type": "feature_card", "title": "数据接入", "body": "统一连接业务库、文件和接口数据", "emphasis": "primary"},
                {"id": "feature_card_2", "type": "feature_card", "title": "智能分析", "body": "自动识别趋势、异常和关键指标", "emphasis": "normal"},
            ],
        }

        items = render_task.extract_items(plan, {})
        cards = render_task.extract_cards(plan, items)

        self.assertIn("三层能力矩阵支撑端到端交付", items)
        self.assertEqual(cards[0]["header"], "数据接入")
        self.assertEqual(cards[0]["emphasis"], "primary")
        self.assertEqual(cards[1]["body"], "自动识别趋势、异常和关键指标")

    def test_chart_and_kpi_components_map_to_generator_params(self):
        chart_plan = {
            "components": [
                {
                    "type": "chart",
                    "data": {
                        "labels": ["Q1", "Q2"],
                        "datasets": [{"name": "收入", "values": [120, 180]}],
                    },
                }
            ]
        }
        kpi_plan = {
            "components": [
                {"type": "kpi_metric", "title": "转化率", "body": "较上月提升", "data": {"value": "38%", "label": "转化率", "delta": "+6pp", "baseline": "较上月"}},
            ]
        }

        self.assertEqual(render_task.chart_data(chart_plan, [])["labels"], ["Q1", "Q2"])
        self.assertEqual(render_task.kpis(kpi_plan, [])[0]["value"], "38%")

    def test_accepted_params_preserves_components_for_keyword_generators(self):
        def component_generator(*, title: str, **kwargs):
            return title, kwargs

        components = [{"type": "feature_card", "title": "山水游览"}]
        params = {
            "title": "桂林体验",
            "components": components,
            "layout_variant": "featured_card_plus_grid",
        }

        accepted = render_task.accepted_params(component_generator, params)

        self.assertEqual(accepted, params)
        self.assertIs(accepted["components"], components)

    def test_component_layout_balances_title_cards_and_long_text(self):
        self.assertEqual(choose_card_columns(3, 11.5, compact=False), 3)
        self.assertEqual(clamp_text("桂林山水" * 20, 12), "桂林山水桂林山水桂林山…")
        self.assertEqual(component_accent({"type": "recommendation"}), "accent")
        self.assertEqual(component_accent({"type": "risk_item"}), "secondary")

    def test_component_label_keeps_paragraph_body_separate(self):
        item = {"type": "paragraph", "text": "桂林的吸引力不只来自单个景点，而是山、水、城共同组织出的连续体验。"}

        self.assertEqual(component_label(item, 1), "要点")
        self.assertTrue(component_body(item).startswith("桂林的吸引力"))

    def test_argument_block_and_list_render_as_narrative_panel(self):
        long_body = (
            "规划质量不能只依赖页数和模板名称，因为真正影响观感的是每页是否有明确论点、证据和结论。"
            "当页面只包含泛化短句时，生成器即使排版稳定，也只能产出看起来正确但内容空泛的页面。"
            "argument_block 允许 Planner 把一段完整论述交给生成器，由生成器负责放入可读正文面板。"
        )
        items = ["先判断本页主张是否具体", "再补充证据、案例或比较口径", "最后给出可以上屏的结论"]
        self.assertEqual(component_items({"type": "list", "items": items}), items)

        prs = render_component_slide(
            palette="ocean_soft",
            title="规划为什么需要长论述",
            content_type="content_slide",
            components=[
                {"type": "argument_block", "title": "核心判断", "body": long_body},
                {"type": "list", "title": "检查顺序", "items": items},
                {"type": "insight", "title": "结果", "body": "长论述和列表拆开后，页面内容更接近可交付的汇报页。"},
            ],
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "narrative-list.pptx"
            save_slide(prs.slides[0], str(output))

            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn("argument_block 允许 Planner", slide_xml)
        self.assertIn("先判断本页主张是否具体", slide_xml)
        self.assertIn("长论述和列表拆开后", slide_xml)

    def test_list_only_narrative_panel_uses_available_space(self):
        items = [
            "自主可控不只是国产替代口号，而是围绕高端芯片、基础软件和关键材料建立可验证供应链。",
            "智能制造要从单点设备升级转向生产数据、质量追溯和柔性排产联动，形成可复制的工厂能力。",
            "绿色低碳需要把能耗、排放和循环利用纳入制造过程指标，让效率提升和减排目标同步落地。",
            "开放合作的重点从规模出口转向标准、品牌和服务体系出海，降低产业链外部波动影响。",
        ]

        prs = render_component_slide(
            palette="ocean_soft",
            title="十五五规划与制造强国新征程",
            subtitle="四个方向都需要从纲要词扩展为可执行判断",
            content_type="content_slide",
            components=[
                {"type": "numbered_list", "title": "重点方向", "items": items},
            ],
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "list-only-panel.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn(items[-1], slide_xml)
        self.assertIn('sz="1340"', slide_xml)

    def test_argument_block_keeps_long_text_without_early_clamp(self):
        long_body = (
            "低空经济并不是单一飞行器赛道，而是由空域管理、飞行服务、制造供应链、运营平台和多行业应用共同构成的新型产业系统。"
            "在政策侧，低空空域改革、试点城市建设和适航审定机制正在同步推进，决定了企业不能只看产品性能，还要判断航线审批、运行监管和应急保障是否成熟。"
            "在商业侧，巡检、物流、文旅、应急和城市治理的需求强度不同，收入模型也从一次性设备销售转向持续运营、数据服务和场景托管。"
            "因此，规划页需要保留完整论述，把概念边界、政策变量、产业链位置和落地条件串联起来，最终形成可执行的场景选择和投资判断。"
        )

        prs = render_component_slide(
            palette="ocean_soft",
            title="低空经济：概念、政策与产业定位",
            content_type="content_slide",
            components=[
                {"type": "argument_block", "title": "核心判断", "body": long_body},
            ],
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "long-argument.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn("最终形成可执行的场景选择和投资判断", slide_xml)

    def test_image_text_embeds_local_image_component(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            image_path = tmp_path / "assets" / "images" / "scene.jpg"
            image_path.parent.mkdir(parents=True)
            Image.new("RGB", (640, 360), color=(40, 120, 180)).save(image_path)

            prs = render_component_slide(
                palette="ocean_soft",
                title="城市低空配送场景",
                content_type="image_text",
                components=[
                    {
                        "type": "image",
                        "title": "无人机配送",
                        "local_path": str(image_path),
                        "caption": "Photo by Demo on Unsplash",
                    },
                    {
                        "type": "argument_block",
                        "title": "场景判断",
                        "body": "城市低空配送需要同时满足可见航线、稳定起降点和末端交付网络三个条件，图片用于呈现真实运行场景，文字负责解释业务约束和落地条件。",
                    },
                ],
            )

            output = tmp_path / "image-text.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                names = package.namelist()
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertTrue(any(name.startswith("ppt/media/") for name in names), names)
        self.assertIn("Photo by Demo on Unsplash", slide_xml)

    def test_image_text_rejects_legacy_asset_id(self):
        with self.assertRaisesRegex(ValueError, "legacy asset id is unsupported"):
            render_component_slide(
                palette="ocean_soft",
                title="旧素材 ID 应直接暴露",
                content_type="image_text",
                components=[
                    {"type": "image", "local_path": "asset:photo_business_work"},
                    {"type": "paragraph", "body": "旧 asset id 不再解析为离线素材。"},
                ],
            )

    def test_image_text_rejects_missing_local_image_path(self):
        with self.assertRaisesRegex(FileNotFoundError, "image path does not exist"):
            render_component_slide(
                palette="ocean_soft",
                title="缺失图片路径应失败",
                content_type="image_text",
                components=[
                    {"type": "image", "local_path": "missing-image.jpg"},
                    {"type": "paragraph", "body": "显式图片字段必须指向真实本地文件。"},
                ],
            )

    def test_image_text_layout_variants_change_regions(self):
        left_image, left_text, _ = image_text_regions("image_left", 0.65, 1.72, 11.95, 5.0)
        right_image, right_text, _ = image_text_regions("image_right", 0.65, 1.72, 11.95, 5.0)
        top_image, top_text, _ = image_text_regions("image_top_band", 0.65, 1.72, 11.95, 5.0)
        bottom_image, bottom_text, _ = image_text_regions("image_bottom_band", 0.65, 1.72, 11.95, 5.0)

        self.assertLess(left_image[0], left_text[0])
        self.assertGreater(right_image[0], right_text[0])
        self.assertEqual(top_image[0], 0.65)
        self.assertGreater(top_text[1], top_image[1] + top_image[3])
        self.assertLessEqual(top_text[1] + top_text[3], 1.72 + 5.0)
        self.assertGreater(bottom_image[1], bottom_text[1])
        self.assertLessEqual(bottom_image[1] + bottom_image[3], 1.72 + 5.0)

    def test_image_text_empty_variant_rotates_by_page_index(self):
        variants = [
            render_task.build_params("image_text", {"page_index": page, "title": f"图文页{page}", "content_plan": {"components": []}}, {"tasks": []})["layout_variant"]
            for page in [1, 2, 3, 4, 5]
        ]

        self.assertEqual(variants, ["image_left", "image_right", "image_top_band", "image_bottom_band", "image_left"])

    def test_build_params_resolves_relative_image_path_to_work_dir(self):
        with tempfile.TemporaryDirectory() as tmp:
            work_dir = Path(tmp)
            task = {
                "title": "图文页",
                "content_type": "image_text",
                "content_plan": {
                    "summary": "说明",
                    "components": [
                        {"type": "image", "local_path": "assets/images/scene.jpg", "caption": "本地图片"},
                    ],
                },
            }

            params = render_task.build_params("image_text", task, {"tasks": [task]}, work_dir=work_dir)

        self.assertEqual(params["components"][0]["local_path"], str((work_dir / "assets/images/scene.jpg").resolve()))

    def test_visual_intent_background_local_path_is_promoted_to_slide_background(self):
        with tempfile.TemporaryDirectory() as tmp:
            work_dir = Path(tmp)
            image_path = work_dir / "assets" / "images" / "background.jpg"
            image_path.parent.mkdir(parents=True)
            Image.new("RGB", (960, 540), color=(90, 120, 150)).save(image_path)
            task = {
                "title": "政策环境",
                "content_type": "content_slide",
                "content_plan": {
                    "visual_intent": {
                        "role": "hero_photo",
                        "asset_purpose": "background",
                        "local_path": "assets/images/background.jpg",
                    },
                    "components": [
                        {"type": "argument_block", "title": "判断", "body": "背景图负责建立政策发布场景，正文负责解释关键政策变量。"},
                    ],
                },
            }

            self.assertEqual(render_task.background_from_task(task, work_dir), str(image_path.resolve()))
            params = render_task.build_params("content_slide", task, {"tasks": [task]}, work_dir=work_dir)
            prs = render_component_slide(
                palette="ocean_soft",
                title=params["title"],
                content_type="content_slide",
                background=render_task.background_from_task(task, work_dir),
                components=params["components"],
            )
            alpha_values = [
                node.get("val")
                for node in prs.slides[0].shapes[-1].element.iter()
                if node.tag.endswith("alpha")
            ]
            self.assertEqual(alpha_values, ["392"])
            output = work_dir / "background-slide.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                names = package.namelist()
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertTrue(any(name.startswith("ppt/media/") for name in names), names)
        self.assertIn("背景图负责建立政策发布场景", slide_xml)
        self.assertNotIn("visual_intent_image", slide_xml)

    def test_background_contain_preserves_full_portrait_image_and_canvas_bounds(self):
        with tempfile.TemporaryDirectory() as tmp:
            work_dir = Path(tmp)
            image_path = work_dir / "portrait.png"
            portrait = Image.new("RGB", (300, 900), color=(30, 180, 60))
            portrait.paste((230, 30, 30), (0, 0, 300, 180))
            portrait.paste((30, 60, 230), (0, 720, 300, 900))
            portrait.save(image_path)

            prs = new_presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_image_background(slide, str(image_path), brightness=1.0, fit_mode="contain")
            output = work_dir / "portrait-background.pptx"
            save_presentation(prs, str(output))

            anchors = [shape for shape in slide.shapes if shape.name == "Background image anchor"]
            self.assertEqual(len(anchors), 1)
            self.assertEqual(len(slide.shapes), 1)
            self.assertEqual(anchors[0].left, 0)
            self.assertEqual(anchors[0].top, 0)
            self.assertEqual(anchors[0].width, prs.slide_width)
            self.assertEqual(anchors[0].height, prs.slide_height)

            with zipfile.ZipFile(output) as package:
                media = [name for name in package.namelist() if name.startswith("ppt/media/")]
                self.assertEqual(len(media), 1)
                rendered_background = Image.open(BytesIO(package.read(media[0]))).convert("RGB")

            self.assertEqual(rendered_background.size, (1920, 1080))
            top_center = rendered_background.getpixel((960, 8))
            bottom_center = rendered_background.getpixel((960, 1071))
            self.assertGreater(top_center[0], 180, top_center)
            self.assertGreater(bottom_center[2], 180, bottom_center)

    def test_background_image_palette_keeps_text_readable_and_moves_image_color_to_fills(self):
        image = Image.new("RGB", (960, 540), color=(214, 178, 80))
        colors = background_image_palette(image, palette="ocean_soft")

        self.assertEqual(colors["primary"], PALETTES["ocean_soft"]["primary"])
        self.assertEqual(colors["accent"], PALETTES["ocean_soft"]["primary"])
        self.assertEqual(colors["secondary"], "51616D")
        self.assertNotEqual(colors["primary_fill"], PALETTES["ocean_soft"]["primary"])
        self.assertGreater(int(colors["text"][:2], 16), 0)
        self.assertTrue(colors["light_bg"].startswith(("E", "F")), colors["light_bg"])

    def test_background_image_palette_returns_fill_tokens_for_shapes(self):
        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "background.jpg"
            Image.new("RGB", (960, 540), color=(70, 125, 150)).save(image_path)

            prs = new_presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            colors = set_image_background(slide, str(image_path), palette="government_red")

        self.assertIsInstance(colors, dict)
        self.assertEqual(colors["primary"], PALETTES["government_red"]["primary"])
        self.assertNotEqual(colors["primary_fill"], PALETTES["government_red"]["primary"])
        self.assertNotEqual(colors["accent_fill"], PALETTES["government_red"]["accent"])
        self.assertTrue(colors["light_bg"].startswith(("E", "F")), colors["light_bg"])

    def test_background_runtime_palette_applies_image_color_to_shapes_not_text(self):
        with tempfile.TemporaryDirectory() as tmp:
            work_dir = Path(tmp)
            image_path = work_dir / "yellow-background.jpg"
            Image.new("RGB", (960, 540), color=(214, 178, 80)).save(image_path)

            prs = render_component_slide(
                title="背景取色",
                subtitle="说明文字应保持可读",
                kicker="图文解读",
                content_type="content_slide",
                background=str(image_path),
                components=[{"type": "paragraph", "title": "关键判断", "body": "形状色跟随背景，文字色保持深色。"}],
            )
            output = work_dir / "runtime-palette.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        xml_colors = set(re.findall(r'val="([0-9A-F]{6})"', slide_xml))
        static_colors = set(PALETTES["ocean_soft"].values())
        readable_text_colors = {"17202A", "51616D", "5A8AA8", "FFFFFF"}
        derived_shape_colors = xml_colors - static_colors - readable_text_colors
        self.assertTrue(derived_shape_colors, xml_colors)
        self.assertIn("51616D", xml_colors)
        self.assertIn("17202A", xml_colors)

    def test_all_background_slides_receive_blur(self):
        with tempfile.TemporaryDirectory() as tmp:
            image_path = Path(tmp) / "background.jpg"
            Image.new("RGB", (640, 360), color=(90, 120, 150)).save(image_path)

            with patch(
                "generators.component_layout.set_image_background",
                return_value=PALETTES["ocean_soft"],
            ) as background_mock:
                render_component_slide(
                    title="普通内容页",
                    content_type="content_slide",
                    background=str(image_path),
                    components=[{"type": "key_point", "body": "正文页背景也需要轻度模糊。"}],
                )
                render_component_slide(
                    title="章节页",
                    content_type="section_divider",
                    background=str(image_path),
                    components=[{"type": "section_marker", "text": "01"}],
                )

            self.assertEqual(background_mock.call_args_list[0].kwargs["blur_radius"], 2)
            self.assertEqual(background_mock.call_args_list[1].kwargs["blur_radius"], 4)

    def test_agenda_uses_manifest_titles_not_summary_blob(self):
        manifest = {
            "tasks": [
                {"task_id": "1", "page_index": 1, "content_type": "title_slide", "title": "低空经济"},
                {
                    "task_id": "2",
                    "page_index": 2,
                    "content_type": "agenda",
                    "title": "目录",
                    "content_plan": {"summary": "第一章 概念与政策 / 第二章 市场规模 / 第三章 城市场景"},
                },
                {"task_id": "3", "page_index": 3, "content_type": "section_divider", "title": "概念、政策与产业定位"},
                {"task_id": "4", "page_index": 4, "content_type": "section_divider", "title": "市场规模与增长动能"},
                {"task_id": "5", "page_index": 5, "content_type": "section_divider", "title": "城市级应用场景全景"},
            ]
        }

        params = render_task.build_params("agenda", manifest["tasks"][1], manifest)

        self.assertEqual(params["items"], ["01  概念、政策与产业定位", "02  市场规模与增长动能", "03  城市级应用场景全景"])
        self.assertNotIn("/", " ".join(params["items"]))
        self.assertEqual([c["type"] for c in params["components"]], ["toc_item", "toc_item", "toc_item"])

    def test_agenda_splits_compact_slash_joined_toc_item(self):
        prs = render_component_slide(
            palette="ocean_soft",
            title="目录",
            content_type="agenda",
            components=[
                {
                    "type": "toc_item",
                    "title": "目录",
                    "body": "第一章 概念与政策 / 第二章 市场规模 / 第三章 城市场景 / 第四章 商业化路径 / 第五章 风险与建议",
                }
            ],
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "agenda.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn("第一章 概念与政策", slide_xml)
        self.assertIn("第五章 风险与建议", slide_xml)
        self.assertNotIn("围绕主题展开结构化说明", slide_xml)

    def test_explicit_components_do_not_duplicate_kpis(self):
        explicit = [
            {
                "type": "kpi_metric",
                "title": "转化率",
                "body": "较上月提升 6 个百分点",
                "data": {"value": "38%", "label": "转化率", "delta": "+6pp", "baseline": "较上月"},
            }
        ]
        prs = generate_kpi_dashboard(
            title="经营指标",
            kpis=[{"value": "38%", "label": "转化率", "baseline": "较上月"}],
            components=explicit,
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "kpi.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertEqual(slide_xml.count("38%"), 1)
        self.assertEqual(slide_xml.count("转化率"), 1)

    def test_comparison_table_renders_structured_table(self):
        prs = generate_comparison_table(
            title="方案选型",
            headers=["旧维度", "旧方案"],
            rows=[["旧数据", "旧结论"]],
            recommendation="旧建议",
            components=[
                {
                    "type": "comparison_matrix",
                    "data": {
                        "headers": ["维度", "方案A", "方案B"],
                        "rows": [
                            ["上线成本", "低，复用现有生成器", "高，需要重写渲染链路"],
                            ["内容质量", "组件先行，便于审查", "页面先行，容易泛化"],
                        ],
                        "highlight_column": "方案A",
                    },
                },
                {"type": "recommendation", "body": "优先选择方案A，先修复组件契约和表格渲染。"},
            ],
        )

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "comparison.pptx"
            save_slide(prs.slides[0], str(output))
            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn("方案A", slide_xml)
        self.assertIn("上线成本", slide_xml)
        self.assertIn("优先选择方案A", slide_xml)
        self.assertNotIn("旧数据", slide_xml)
        self.assertNotIn("围绕主题展开结构化说明", slide_xml)

    def test_source_footer_compacts_long_urls(self):
        source = "来源: 国家统计局 2025年数据 | https://www.stats.gov.cn/sj/zxfb/202502/t20250228_1958901.html；World Bank | https://data.worldbank.org/indicator/NY.GDP.MKTP.CD"
        compact = compact_source_text(source, max_chars=60)

        self.assertLessEqual(len(compact), 60)
        self.assertIn("含2个链接", compact)
        self.assertNotIn("https://", compact)

    def test_sourced_deep_dive_panels_stay_above_footer(self):
        prs = render_component_slide(
            palette="ocean_soft",
            title="AI落地的三大挑战与应对",
            subtitle="深入分析信任成本、系统交付和人才能力",
            source="来源: Gartner 2026 AI Survey、毕马威AI投资回报报告",
            content_type="deep_dive",
            components=[
                {
                    "type": "argument_block",
                    "body": "AI落地正面临高热度与窄落地的结构性反差。" * 28,
                },
                {"type": "risk_item", "body": "工业场景的信任成本高。"},
                {"type": "risk_item", "body": "系统交付能力仍然不足。"},
                {"type": "opportunity_item", "body": "协议标准化正在改善工具协作。"},
                {"type": "recommendation", "body": "通过分层培训补齐复合型人才能力。"},
            ],
        )

        slide = prs.slides[0]
        footer_divider = next(
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and abs(shape.top - Inches(6.85)) <= 1
            and abs(shape.height - Inches(0.01)) <= 1
        )
        content_panels = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and abs(shape.top - Inches(1.72)) <= 1
            and shape.height > Inches(4.0)
        ]

        self.assertEqual(len(content_panels), 1)
        for panel in content_panels:
            self.assertLessEqual(panel.top + panel.height, footer_divider.top - Inches(0.2))

    def test_atomic_components_render_without_coordinates(self):
        prs = render_component_slide(
            palette="charcoal_light",
            title="组件式执行链路",
            content_type="deep_dive",
            components=[
                {"type": "tag", "text": "准备阶段"},
                {"type": "icon", "icon": "LLM", "title": "规划器"},
                {"type": "divider", "role": "区分规划与渲染"},
                {"type": "shape", "role": "强调组件计划已锁定"},
                {"type": "architecture_box", "title": "Deck Planner", "body": "负责整套叙事、章节和页面角色，不接触坐标、字号或颜色。", "role": "规划层"},
                {"type": "architecture_box", "title": "Component Planner", "body": "把每页拆成可执行语义组件，供生成器稳定消费。", "role": "规划层"},
                {"type": "arrow", "relation": "输出 DeckSpec", "target": "box_2"},
            ],
        )

        self.assertEqual(len(prs.slides), 1)

    def test_saved_slide_includes_default_transition(self):
        prs = new_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "transition.pptx"
            save_slide(slide, str(output))

            with zipfile.ZipFile(output) as package:
                slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertIn("<p:transition", slide_xml)
        self.assertIn("<p:fade", slide_xml)
        self.assertIn('advClick="1"', slide_xml)


if __name__ == "__main__":
    unittest.main()
