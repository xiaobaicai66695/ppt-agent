"""PPT generator base module - shared utilities and palette definitions."""
from __future__ import annotations

import os
import re
from copy import deepcopy
from typing import Literal

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Palette definitions (浅色基调，饱和度 5%-45%，亮度 80%-95%)
# ---------------------------------------------------------------------------

PALETTES: dict[str, dict[str, str]] = {
    "ocean_soft": {
        "name": "雾霾蓝",
        "primary": "5A8AA8",
        "secondary": "7BA3B8",
        "accent": "A8C4D4",
        "background": "F5F8FA",
        "text": "3A4A52",
        "text_muted": "7A8A92",
        "light_bg": "E8F0F5",
        "divider": "C8D8E5",
    },
    "sage_calm": {
        "name": "鼠尾草绿",
        "primary": "6A8E85",
        "secondary": "84B59F",
        "accent": "A8C5BD",
        "background": "FAF5F2",
        "text": "3A4A46",
        "text_muted": "8A9A92",
        "light_bg": "E8F0EC",
        "divider": "C8D8D0",
    },
    "warm_terracotta": {
        "name": "陶土橙",
        "primary": "C47060",
        "secondary": "D4A574",
        "accent": "E8DDD0",
        "background": "FAF5F2",
        "text": "4A3A32",
        "text_muted": "9A8A82",
        "light_bg": "F0E8E0",
        "divider": "D8C8B8",
    },
    "charcoal_light": {
        "name": "浅炭灰",
        "primary": "5A6A75",
        "secondary": "8A9AA5",
        "accent": "C8D0D5",
        "background": "F8F8F8",
        "text": "3A4248",
        "text_muted": "8A9298",
        "light_bg": "E8ECEE",
        "divider": "C0C8CC",
    },
    "berry_cream": {
        "name": "玫瑰灰粉",
        "primary": "8D5A6B",
        "secondary": "B89098",
        "accent": "E0D0D4",
        "background": "FAF5F5",
        "text": "4A3238",
        "text_muted": "9A8288",
        "light_bg": "F0E4E8",
        "divider": "D8C0C8",
    },
    "lavender_mist": {
        "name": "薰衣草灰",
        "primary": "7A6A8A",
        "secondary": "9A8AAA",
        "accent": "C8B8D4",
        "background": "F8F5FA",
        "text": "3A3248",
        "text_muted": "8A7A98",
        "light_bg": "EDE4F4",
        "divider": "D0C0E0",
    },
    "forest_moss": {
        "name": "苔藓绿",
        "primary": "5A8E5A",
        "secondary": "7EBF7E",
        "accent": "A8D4A8",
        "background": "F5FAF5",
        "text": "1A2A1A",
        "text_muted": "5A7A5A",
        "light_bg": "E4F4E4",
        "divider": "A8C8A8",
    },
    "sunset_peach": {
        "name": "杏色",
        "primary": "C4A080",
        "secondary": "D8B898",
        "accent": "E8D8C8",
        "background": "FAF5F2",
        "text": "3A2010",
        "text_muted": "8A6A4A",
        "light_bg": "F0E4D8",
        "divider": "D8C8B0",
    },

    # ── 中国场景配色 ──
    "government_red": {
        "name": "政务红",
        "primary": "8B1A1A",
        "secondary": "B87070",
        "accent": "D4A0A0",
        "background": "FDF8F5",
        "text": "3A1A1A",
        "text_muted": "8A6A6A",
        "light_bg": "F5E8E5",
        "divider": "D8C0B8",
    },
    "patriotic_blue": {
        "name": "爱国蓝",
        "primary": "2B5797",
        "secondary": "4A7AB0",
        "accent": "A0B8D4",
        "background": "F5F8FC",
        "text": "1A2A48",
        "text_muted": "5A6A8A",
        "light_bg": "E4EAF4",
        "divider": "B8C8D8",
    },
    "civic_gold": {
        "name": "公民金",
        "primary": "B8860B",
        "secondary": "C8A040",
        "accent": "E8D890",
        "background": "FAF8F0",
        "text": "3A2A0A",
        "text_muted": "8A7A4A",
        "light_bg": "F0E8D8",
        "divider": "D8C8A0",
    },
    "debate_purple": {
        "name": "答辩紫",
        "primary": "5A4A7A",
        "secondary": "8A74AA",
        "accent": "C8B8D4",
        "background": "F8F5FA",
        "text": "2A1A3A",
        "text_muted": "7A6A8A",
        "light_bg": "EDE4F4",
        "divider": "C8B8D8",
    },
    "activity_orange": {
        "name": "活力橙",
        "primary": "C06030",
        "secondary": "C89060",
        "accent": "E8C8A0",
        "background": "FAF5F2",
        "text": "3A2010",
        "text_muted": "8A6A4A",
        "light_bg": "F0E0D0",
        "divider": "D8C0A8",
    },
    "report_green": {
        "name": "报告绿",
        "primary": "3A6A4A",
        "secondary": "5A8A6A",
        "accent": "A0C0B0",
        "background": "F5F8F5",
        "text": "1A2A1A",
        "text_muted": "5A7A5A",
        "light_bg": "E4F0E8",
        "divider": "A8C0B0",
    },
    "simple_gray": {
        "name": "简约灰",
        "primary": "4A4A4A",
        "secondary": "6A6A6A",
        "accent": "C0C0C0",
        "background": "F8F8F8",
        "text": "2A2A2A",
        "text_muted": "8A8A8A",
        "light_bg": "E8E8E8",
        "divider": "C0C0C0",
    },
}


def rgb(hex_color: str) -> RGBColor:
    """Convert hex string like '5A8AA8' to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# ---------------------------------------------------------------------------
# Presentation helpers
# ---------------------------------------------------------------------------

def new_presentation(
    width_in: float = 13.333,
    height_in: float = 7.5,
    palette: str = "ocean_soft",
) -> Presentation:
    """Create a new 16:9 presentation. No slide is created — generators add slides."""
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def set_slide_background(slide, palette: str = "ocean_soft"):
    """Set solid background color for a slide."""
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(colors["background"])


def set_image_background(
    slide,
    image_path: str,
    brightness: float = 0.95,
    palette: str = "ocean_soft",
    blur_radius: float = 0,
    fit_mode: Literal["contain", "cover"] = "cover",
):
    """
    为幻灯片设置图片背景。

    Args:
        slide: python-pptx Slide 对象
        image_path: 背景图片的完整路径（支持 JPG/PNG）
        brightness: 亮度调整 (0.0-1.0)，值越小背景越暗。推荐值 0.9-1.0
        blur_radius: 高斯模糊半径。标题页和章节分割页使用轻度模糊，
            避免复杂背景或透视结构压低文字可读性。
        fit_mode: 背景适配方式。默认 cover，等比铺满画布并允许边缘适度裁剪；
            contain 完整保留原图，并用同图的模糊扩展层填满比例不一致的留白。

    Example:
        set_image_background(slide, "D:/path/to/background.jpg")
    """
    if not os.path.exists(image_path):
        print(f"Warning: Background image not found: {image_path}")
        return

    with Image.open(image_path) as source:
        img = source.convert("RGB")

    if brightness < 1.0:
        img = ImageEnhance.Brightness(img).enhance(brightness)

    # 按当前演示文稿的真实宽高比生成背景，不再假设所有页面都固定为 16:9。
    presentation = slide.part.package.presentation_part.presentation
    slide_w, slide_h = presentation.slide_width, presentation.slide_height
    if slide_w >= slide_h:
        target_w = 1920
        target_h = max(1, round(target_w * slide_h / slide_w))
    else:
        target_h = 1920
        target_w = max(1, round(target_h * slide_w / slide_h))

    img = _fit_background_image(img, (target_w, target_h), fit_mode)
    if blur_radius > 0:
        img = img.filter(ImageFilter.GaussianBlur(radius=blur_radius))

    # 临时保存处理后的图片
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        temp_path = tmp.name
    img.save(temp_path, quality=95)

    # 保留一个与画布严格同尺寸的底层图片锚点，确保 python-pptx 保存/拆页后
    # 媒体 relationship 不会丢失，同时避免原图原生尺寸造成整页越界。
    img_shape = slide.shapes.add_picture(
        temp_path,
        0,
        0,
        width=slide_w,
        height=slide_h,
    )
    img_shape.name = "Background image anchor"
    # 只保留这一个底层图片来源。原生 p:bg 会随测试时扩大的页面继续铺满，
    # 既造成越界误报，也会让部分查看器和图片锚点出现双重缩放差异。
    sp_tree = slide.shapes._spTree
    sp_tree.remove(img_shape._element)
    sp_tree.insert(2, img_shape._element)

    os.remove(temp_path)

    # 自动添加磨砂玻璃叠加层，提升内容可读性
    # add_frosted_glass_overlay 返回适合玻璃背景的颜色映射
    return add_frosted_glass_overlay(slide, palette)


def _fit_background_image(
    image: Image.Image,
    target_size: tuple[int, int],
    fit_mode: Literal["contain", "cover"] = "cover",
) -> Image.Image:
    """Return a slide-sized background without stretching the source image."""
    if fit_mode == "cover":
        return ImageOps.fit(
            image,
            target_size,
            method=Image.Resampling.LANCZOS,
            centering=(0.5, 0.5),
        )
    if fit_mode != "contain":
        raise ValueError(f"Unsupported background fit_mode: {fit_mode}")

    # 先生成一层铺满画布的模糊同图背景，再将完整原图等比居中覆盖。
    # 这样既能保留全部画面，也不会出现黑边、白边或非等比拉伸。
    backdrop = ImageOps.fit(
        image,
        target_size,
        method=Image.Resampling.LANCZOS,
        centering=(0.5, 0.5),
    )
    backdrop = ImageEnhance.Brightness(backdrop).enhance(0.82)
    backdrop = backdrop.filter(
        ImageFilter.GaussianBlur(radius=max(18, round(min(target_size) * 0.025)))
    )

    foreground = ImageOps.contain(
        image,
        target_size,
        method=Image.Resampling.LANCZOS,
    )
    left = (target_size[0] - foreground.width) // 2
    top = (target_size[1] - foreground.height) // 2
    backdrop.paste(foreground, (left, top))
    return backdrop


def add_frosted_glass_overlay(slide, palette: str = "ocean_soft"):
    """
    在幻灯片上添加磨砂玻璃叠加层并返回适合的颜色映射。

    磨砂玻璃覆盖整张幻灯片，顶层所有文字和色块印在上面。
    返回的颜色映射中 accent/divider 等替换为适合玻璃背景的白色/深色，
    避免和背景图片的蓝色调违和。

    Returns:
        dict: 颜色映射，key 为 palette 颜色名，value 为适合玻璃背景的新颜色。
              例如 {"primary": "FFFFFF", "divider": "D0D0D0", "accent": "FFFFFF"}
    """
    add_rect(
        slide,
        left=0.0, top=0.0, width=13.333, height=7.5,
        fill_color=(255, 255, 255, 96),
        line_color=None,
    )
    # 玻璃模式下，文字和色块用深色系（在白色磨砂底上可读），
    # 填充色用深蓝色系（保持海洋配色风格，在磨砂玻璃上有足够对比度）
    glass_palette = PALETTES.get(palette, PALETTES["ocean_soft"]).copy()
    glass_palette.update({
        # 文字色 → 深色（在白色磨砂底上可读）
        "primary": "2D3748",   # 深灰（标题/小标题）
        "secondary": "4A5568",  # 中灰（小标题/kicker）
        "accent": "2C5282",     # 深蓝（强调）
        "text": "1A202C",      # 近黑（正文）
        "text_muted": "718096", # 中灰（次要文字）
        # 填充色 → 深蓝（保持海洋风格又有对比度）
        "divider": "A0B4C0",
        "light_bg": "C8D8E5",
        "background": "D8E8F0",
        # 象限色（SWOT等用原始色块时，深蓝比浅蓝对比度更好）
        "primary_fill": "3D6B82",
        "secondary_fill": "5A8AA8",
        "accent_fill": "7BA3B8",
    })
    return glass_palette


def default_slide_transition() -> str:
    """Return the default slide transition effect for exported PPTX files."""
    return os.environ.get("PPT_SLIDE_TRANSITION", "fade").strip().lower()


def default_slide_transition_speed() -> str:
    """Return the default transition speed accepted by PowerPoint OOXML."""
    speed = os.environ.get("PPT_SLIDE_TRANSITION_SPEED", "med").strip().lower()
    return speed if speed in {"slow", "med", "fast"} else "med"


def _transition_effect_xml(effect: str, direction: str):
    effect = (effect or "fade").strip().lower()
    direction = (direction or "l").strip().lower()
    supported_dirs = {"l", "r", "u", "d", "lu", "ld", "ru", "rd"}
    if direction not in supported_dirs:
        direction = "l"

    if effect == "none":
        return None
    if effect == "push":
        node = etree.Element(qn("p:push"))
        node.set("dir", direction)
        return node
    if effect == "wipe":
        node = etree.Element(qn("p:wipe"))
        node.set("dir", direction)
        return node
    if effect == "split":
        node = etree.Element(qn("p:split"))
        node.set("orient", os.environ.get("PPT_SLIDE_TRANSITION_ORIENT", "horz"))
        node.set("dir", os.environ.get("PPT_SLIDE_TRANSITION_SPLIT_DIR", "out"))
        return node
    if effect == "cover":
        node = etree.Element(qn("p:cover"))
        node.set("dir", direction)
        return node
    if effect == "uncover":
        node = etree.Element(qn("p:uncover"))
        node.set("dir", direction)
        return node
    node = etree.Element(qn("p:fade"))
    return node


def apply_slide_transition(
    slide,
    effect: str | None = None,
    speed: str | None = None,
    direction: str | None = None,
):
    """Apply a PowerPoint slide transition using raw OOXML.

    python-pptx does not expose slide transitions, so the generator writes the
    minimal OOXML PowerPoint expects. The transition is intentionally modest by
    default: click-to-advance fade at medium speed.
    """
    effect = default_slide_transition() if effect is None else effect.strip().lower()
    if effect == "none":
        remove_slide_transition(slide)
        return

    speed = default_slide_transition_speed() if speed is None else speed.strip().lower()
    if speed not in {"slow", "med", "fast"}:
        speed = "med"
    direction = os.environ.get("PPT_SLIDE_TRANSITION_DIRECTION", "l") if direction is None else direction

    transition = etree.Element(qn("p:transition"))
    transition.set("spd", speed)
    transition.set("advClick", "1")

    effect_node = _transition_effect_xml(effect, direction)
    if effect_node is not None:
        transition.append(effect_node)

    slide_el = slide._element
    remove_slide_transition(slide)
    insert_after = None
    for tag in (qn("p:clrMapOvr"), qn("p:cSld")):
        candidate = slide_el.find(tag)
        if candidate is not None:
            insert_after = candidate
            break
    if insert_after is None:
        slide_el.insert(0, transition)
        return
    slide_el.insert(list(slide_el).index(insert_after) + 1, transition)


def remove_slide_transition(slide):
    """Remove existing slide transition XML if present."""
    slide_el = slide._element
    existing = slide_el.find(qn("p:transition"))
    if existing is not None:
        slide_el.remove(existing)


def apply_presentation_transitions(prs: Presentation):
    """Apply default transitions to every slide in a presentation."""
    if default_slide_transition() == "none":
        for slide in prs.slides:
            remove_slide_transition(slide)
        return
    for slide in prs.slides:
        apply_slide_transition(slide)


def save_presentation(prs: Presentation, output_path: str):
    """Save the presentation to a file path."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    apply_presentation_transitions(prs)
    prs.save(output_path)


def _clone_xml(el):
    """Clone an lxml element by round-tripping through string serialization.
    This produces a clean copy detached from the source document tree,
    avoiding namespace conflicts that deepcopy can cause across different Presentations.
    """
    from lxml import etree

    return etree.fromstring(etree.tostring(el, encoding="unicode"))


def save_slide(slide, output_path: str):
    """
    Save a single slide as a standalone PPTX file.
    Creates a new presentation, clones the slide (including chart parts and
    background image resources) into it, and saves.
    """
    from copy import deepcopy
    from lxml import etree
    from pptx.oxml.ns import qn
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    src_slide_part = slide.part
    src_slide_rel = src_slide_part._rels

    new_prs = Presentation()
    new_prs.slide_width = src_slide_part.package.presentation_part.presentation.slide_width
    new_prs.slide_height = src_slide_part.package.presentation_part.presentation.slide_height

    blank_layout = new_prs.slide_layouts[6]
    dest_slide = new_prs.slides.add_slide(blank_layout)
    dest_slide_part = dest_slide.part

    # Remove placeholder shapes from the blank slide
    for shape in list(dest_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # ------------------------------------------------------------------
    # Copy background: handle both solid-color and image (blipFill) backgrounds.
    # ------------------------------------------------------------------
    p_cSld = dest_slide._element.find(qn("p:cSld"))
    # Remove any existing background
    existing_bg = p_cSld.find(qn("p:bg"))
    if existing_bg is not None:
        p_cSld.remove(existing_bg)

    src_p_cSld = slide._element.find(qn("p:cSld"))
    src_bg = src_p_cSld.find(qn("p:bg")) if src_p_cSld is not None else None

    if src_bg is not None:
        bg_clone = deepcopy(src_bg)
        # If it's a blipFill (image) background, copy the image part
        bgPr = bg_clone.find(qn("p:bgPr"))
        if bgPr is not None:
            blipFill = bgPr.find(qn("a:blipFill"))
            if blipFill is not None:
                blip = blipFill.find(qn("a:blip"))
                if blip is not None:
                    old_rid = blip.get(qn("r:embed"))
                    if old_rid and old_rid in src_slide_rel:
                        try:
                            img_part = src_slide_rel[old_rid].target_part
                            new_rid = dest_slide_part.relate_to(img_part, RT.IMAGE)
                            blip.set(qn("r:embed"), new_rid)
                        except Exception:
                            pass
        p_cSld.insert(0, bg_clone)

    # ------------------------------------------------------------------
    # Copy chart relationships.
    # ------------------------------------------------------------------
    for rel_id, rel in src_slide_rel.items():
        if rel.reltype == RT.CHART:
            chart_part = rel.target_part
            dest_slide_part.relate_to(chart_part, RT.CHART)

    # ------------------------------------------------------------------
    # Clone all shape elements.
    # Update chart rIds to point to the newly-added chart parts.
    # ------------------------------------------------------------------
    dest_spTree = dest_slide.shapes._spTree
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for shape in slide.shapes:
        el = deepcopy(shape.element)
        chart_el = el.find(qn("c:chart"))
        if chart_el is not None:
            old_rid = chart_el.get("{" + r_ns + "}embed")
            if old_rid and old_rid in src_slide_rel:
                src_rel = src_slide_rel[old_rid]
                if src_rel.reltype == RT.CHART:
                    for new_rel_id, new_rel in dest_slide_part.rels.items():
                        if new_rel.target_part == src_rel.target_part:
                            chart_el.set("{" + r_ns + "}embed", new_rel_id)
                            break
        for blip in el.findall(".//" + qn("a:blip")):
            old_rid = blip.get("{" + r_ns + "}embed")
            if old_rid and old_rid in src_slide_rel:
                src_rel = src_slide_rel[old_rid]
                if src_rel.reltype == RT.IMAGE:
                    try:
                        img_part = src_rel.target_part
                        new_rid = dest_slide_part.relate_to(img_part, RT.IMAGE)
                        blip.set("{" + r_ns + "}embed", new_rid)
                    except Exception:
                        pass
        dest_spTree.append(el)

    save_presentation(new_prs, output_path)


def compact_source_text(source: str, max_chars: int = 96) -> str:
    """Return a footer-safe source label while avoiding long URL spillover."""
    text = str(source or "").strip()
    if not text:
        return ""
    urls = re.findall(r"https?://[^\s|；;，,）)]+", text)
    without_urls = re.sub(r"https?://[^\s|；;，,）)]+", "", text)
    parts = [
        part.strip(" |；;，,")
        for part in re.split(r"[|；;]+", without_urls)
        if part.strip(" |；;，,")
    ]
    compact = "；".join(parts[:3]) if parts else "来源"
    if urls and "链接" not in compact:
        compact = f"{compact}（含{len(urls)}个链接）"
    if not compact.startswith(("来源", "Source", "source")):
        compact = f"来源: {compact}"
    if len(compact) > max_chars:
        compact = compact[: max_chars - 1].rstrip(" |；;，,") + "…"
    return compact


def add_source_line(
    slide,
    source: str,
    palette: str = "ocean_soft",
):
    """在幻灯片底部渲染数据来源/参考信息。
    当 source 非空时，在底部添加灰色小字来源行。
    """
    source = compact_source_text(source)
    if not source:
        return

    # 底部淡灰色分割线
    add_rect(
        slide,
        left=0.5, top=6.85, width=12.0, height=0.01,
        fill_color="text_muted", palette=palette,
    )

    # 来源文字（小号灰色）
    add_text(
        slide,
        text=source,
        left=0.5, top=6.9, width=12.0, height=0.35,
        font_size=9, bold=False,
        color="text_muted", alignment="left",
        palette=palette,
    )


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400


def _text_units(text: str) -> float:
    """Estimate rendered text width; CJK glyphs are wider than ASCII."""
    units = 0.0
    for ch in str(text or ""):
        if "\u4e00" <= ch <= "\u9fff" or "\u3000" <= ch <= "\u303f" or "\uff00" <= ch <= "\uffef":
            units += 1.0
        elif ch.isspace():
            units += 0.35
        else:
            units += 0.55
    return units


def _estimate_lines(text: str, width_in: float, font_size: float) -> int:
    usable_width = max(width_in - 0.08, 0.15)
    units_per_line = max(1.0, usable_width * 72 / max(font_size * 0.92, 1))
    lines = 0
    for raw_line in str(text or "").splitlines() or [""]:
        units = _text_units(raw_line)
        lines += max(1, int((units + units_per_line - 1) // units_per_line))
    return lines


def estimate_text_height_pt(text: str, width_in: float, font_size: float, line_height: float = 1.16) -> float:
    return _estimate_lines(text, width_in, font_size) * font_size * line_height


def fit_font_size(
    text: str,
    width_in: float,
    height_in: float,
    font_size: float,
    min_font_size: float = None,
    max_font_size: float = None,
    line_height: float = 1.16,
    allow_growth: bool = True,
) -> float:
    """Fit text to a box and gently enlarge sparse content when space allows."""
    if min_font_size is None:
        min_font_size = 24 if font_size >= 34 else (12 if font_size >= 14 else font_size)
    if max_font_size is None:
        if font_size >= 30:
            max_font_size = font_size * 1.18
        elif font_size >= 16 and height_in >= 0.45:
            max_font_size = font_size * 1.12
        else:
            max_font_size = font_size

    size = float(font_size)
    while size > min_font_size:
        estimated_height_pt = estimate_text_height_pt(text, width_in, size, line_height)
        if estimated_height_pt <= height_in * 72:
            break
        size -= 1

    if allow_growth and text and max_font_size > size:
        while size + 1 <= max_font_size:
            candidate = size + 1
            estimated_height_pt = estimate_text_height_pt(text, width_in, candidate, line_height)
            if estimated_height_pt > height_in * 72 * 0.82:
                break
            size = candidate

    return max(size, min_font_size)


def balanced_top(region_top: float, region_height: float, content_height: float, min_top: float = None) -> float:
    """Return a top coordinate that centers content within a vertical region."""
    top = region_top + max(0.0, (region_height - content_height) / 2)
    if min_top is not None:
        top = max(min_top, top)
    return top


def _auto_vertical_anchor(text: str, width: float, height: float, font_size: float, alignment: str) -> MSO_ANCHOR:
    lines = _estimate_lines(text, width, font_size)
    text_height = estimate_text_height_pt(text, width, font_size)
    box_height = height * 72
    if box_height >= text_height * 1.28 and (alignment == "center" or lines <= 2 or font_size >= 18):
        return MSO_ANCHOR.MIDDLE
    return MSO_ANCHOR.TOP


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    alignment: Literal["left", "center", "right"] = "left",
    palette: str = "ocean_soft",
    font_name: str = None,
    italic: bool = False,
    colors: dict = None,
    min_font_size: float = None,
    max_font_size: float = None,
    vertical_alignment: Literal["top", "middle", "bottom", "center", "auto"] = "auto",
    line_spacing: float = 1.0,
    margin: float = 0.03,
) -> "pptx.shapes.shapetree.Shape":
    """Add a text box to the slide with consistent styling."""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(max(margin * 0.6, 0.01))
    tf.margin_bottom = Inches(max(margin * 0.6, 0.01))

    p = tf.paragraphs[0]
    p.text = text
    fitted_size = fit_font_size(
        text,
        width,
        height,
        font_size,
        min_font_size=min_font_size,
        max_font_size=max_font_size,
        allow_growth=max_font_size is not False,
    )
    p.font.size = Pt(fitted_size)
    p.font.bold = bold
    p.font.italic = italic
    p.line_spacing = line_spacing

    if colors is not None:
        p.font.color.rgb = rgb(colors.get(color, color))
    else:
        palette_colors = PALETTES.get(palette, PALETTES["ocean_soft"])
        p.font.color.rgb = rgb(palette_colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)

    if font_name:
        p.font.name = font_name

    anchor = vertical_alignment
    if anchor == "auto":
        tf.anchor = _auto_vertical_anchor(text, width, height, fitted_size, alignment)
    else:
        tf.anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "center": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(anchor, MSO_ANCHOR.TOP)

    return txbox


def add_text_boxed(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    alignment: Literal["left", "center", "right"] = "left",
    vertical_alignment: Literal["top", "middle", "bottom", "center", "auto"] = "auto",
    palette: str = "ocean_soft",
    colors: dict = None,
    min_font_size: float = None,
    max_font_size: float = None,
    line_spacing: float = 1.0,
):
    """Add a self-balancing text box. Kept separate for explicit generator usage."""
    return add_text(
        slide,
        text=text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_size=font_size,
        bold=bold,
        color=color,
        alignment=alignment,
        vertical_alignment=vertical_alignment,
        palette=palette,
        colors=colors,
        min_font_size=min_font_size,
        max_font_size=max_font_size,
        line_spacing=line_spacing,
    )


def add_paragraph(
    tf,
    text: str,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    palette: str = "ocean_soft",
    alignment: Literal["left", "center", "right"] = "left",
    space_before: float = 0,
    space_after: float = 0,
    italic: bool = False,
    font_name: str = None,
    colors: dict = None,
    min_font_size: float = None,
    line_spacing: float = 1.0,
):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    try:
        shape = tf._parent
        width = shape.width / EMU_PER_INCH
        height = shape.height / EMU_PER_INCH
        fitted_size = fit_font_size(text, width, height, font_size, min_font_size)
    except Exception:
        fitted_size = font_size
    p.font.size = Pt(fitted_size)
    p.font.bold = bold
    p.font.italic = italic
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if colors is not None:
        p.font.color.rgb = rgb(colors.get(color, color))
    else:
        palette_colors = PALETTES.get(palette, PALETTES["ocean_soft"])
        p.font.color.rgb = rgb(palette_colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)
    if font_name:
        p.font.name = font_name
    return p


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add a filled rectangle.

    Args:
        fill_color: either a palette key (str) like "primary", or an (R, G, B, A) tuple.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    if isinstance(fill_color, tuple) and len(fill_color) == 4:
        r, g, b, a = fill_color
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        # Set alpha transparency directly in XML: 0=opaque, 100000=fully transparent
        from pptx.oxml.ns import qn
        spPr = shape.element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(int(a / 255 * 100000)))
    else:
        fill.solid()
        fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_round_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    if isinstance(fill_color, tuple) and len(fill_color) == 4:
        r, g, b, a = fill_color
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        from pptx.oxml.ns import qn
        spPr = shape.element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(int(a / 255 * 100000)))
    else:
        fill.solid()
        fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_glass_panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    palette: str = "ocean_soft",
    fill_color: str = "background",
    alpha: int = 96,
    line_color: str = "divider",
    line_width: float = 0.5,
) -> "pptx.shapes.shapetree.Shape":
    """Add a semi-transparent content panel over image backgrounds."""
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    base = colors.get(fill_color, fill_color)
    color = rgb(base)
    return add_round_rect(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_color=(color[0], color[1], color[2], alpha),
        palette=palette,
        line_color=line_color,
        line_width=line_width,
    )


def add_ellipse(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    palette: str = "ocean_soft",
    line_color: str = None,
    line_width: float = 0,
) -> "pptx.shapes.shapetree.Shape":
    """Add an ellipse/circle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill = shape.fill
    fill.solid()
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    fill.fore_color.rgb = rgb(colors.get(fill_color, fill_color))

    line = shape.line
    if line_color:
        line.color.rgb = rgb(colors.get(line_color, line_color))
        line.width = Pt(line_width)
    else:
        line.fill.background()

    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "primary",
    width: float = 1.5,
    palette: str = "ocean_soft",
) -> "pptx.shapes.shapetree.Shape":
    """Add a straight line connector."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = rgb(PALETTES.get(palette, PALETTES["ocean_soft"]).get(color, color))
    connector.line.width = Pt(width)
    return connector


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "secondary",
    width: float = 1.5,
    palette: str = "ocean_soft",
) -> "pptx.shapes.shapetree.Shape":
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1), Inches(y1), Inches(x2 - x1), Inches(y2 - y1),
    )
    shape.fill.solid()
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    shape.fill.fore_color.rgb = rgb(colors.get(color, color))
    shape.line.fill.background()
    return shape


def add_text_in_shape(
    shape,
    text: str,
    font_size: float = 14,
    bold: bool = False,
    color: str = "text",
    palette: str = "ocean_soft",
    alignment: Literal["left", "center", "right"] = "left",
    vertical_alignment: str = "top",
    font_name: str = None,
    italic: bool = False,
    min_font_size: float = None,
):
    """Set text inside an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Clear existing paragraphs except first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    p = tf.paragraphs[0]
    p.text = text
    width = shape.width / EMU_PER_INCH
    height = shape.height / EMU_PER_INCH
    p.font.size = Pt(fit_font_size(text, width, height, font_size, min_font_size))
    p.font.bold = bold
    p.font.italic = italic
    colors = PALETTES.get(palette, PALETTES["ocean_soft"])
    p.font.color.rgb = rgb(colors.get(color, color))
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)

    if font_name:
        p.font.name = font_name

    # Vertical alignment
    tf.anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
        "center": MSO_ANCHOR.MIDDLE,
    }.get(vertical_alignment, MSO_ANCHOR.TOP)


# ---------------------------------------------------------------------------
# Background utilities
# ---------------------------------------------------------------------------

def get_background_path(
    theme: str = None,
    scenario: str = None,
    random_select: bool = False,
) -> str:
    """Compatibility no-op for the removed local background catalog."""
    from .background_manager import get_background as _get_bg
    return _get_bg(theme=theme, scenario=scenario, random_select=random_select)


def resolve_background(background: str) -> str:
    """Resolve only explicit image file paths.

    Legacy theme/scenario names are no longer resolved; new plans pass
    downloaded image paths through visual_intent.local_path or image components.
    """
    if not background:
        return None

    # 已经是绝对路径
    if os.path.isabs(background) and os.path.exists(background):
        return background

    # 文件存在
    if os.path.exists(background):
        return background

    return None
