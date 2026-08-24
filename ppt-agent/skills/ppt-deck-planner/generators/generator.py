#!/usr/bin/env python3
"""
PPT Generator - Main entry point.

Usage:
    python generator.py

生成一个完整的PPT示例，展示所有29种页面模板类型。
"""
from pathlib import Path
import sys
import os

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from generators import (
    PALETTES,
    new_presentation,
    generate_title_slide,
    generate_section_divider,
    generate_content_slide,
    generate_stat_slide,
    generate_quote_slide,
    generate_card_grid,
    generate_process_flow,
    generate_two_column,
    generate_three_column,
    generate_summary_slide,
    generate_image_text,
    generate_example_detail,
    generate_deep_dive,
    generate_agenda,
    generate_case_study,
    generate_kpi_dashboard,
    generate_comparison_table,
    generate_image_hero,
    generate_kanban,
    generate_brand_focus,
    generate_region_map,
    generate_timeline,
    create_bar_chart_slide,
    create_line_chart_slide,
    create_pie_chart_slide,
    create_table_slide,
    create_handdrawn_style_slide,
    create_icon_combo_slide,
    create_smart_layout_slide,
)
from generators.base import (
    add_source_line, add_text, add_rect, add_ellipse, add_round_rect,
    add_line, set_slide_background, save_presentation, save_slide,
    rgb, PALETTES as BASE_PALETTES,
)


# ============================================================================
# 图表类模板 (需要特殊绘制，无法通过通用模板生成)
# ============================================================================

def create_bar_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """柱状图 - 季度销售数据"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="季度销售业绩",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="各产品线季度收入对比（单位：万元）",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1 一季度", "Q2 二季度", "Q3 三季度", "Q4 四季度"]
    chart_data.add_series("产品A线", (125, 148, 132, 189))
    chart_data.add_series("产品B线", (98, 115, 142, 167))
    chart_data.add_series("产品C线", (67, 78, 95, 112))

    x, y, cx, cy = Inches(0.8), Inches(1.5), Inches(11.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.bold = True

    # Reduce Y-axis tick density
    val_axis = chart.value_axis
    val_axis.major_unit = 50.0
    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 200.0

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)

    return slide


def create_line_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """折线图 - 用户增长趋势"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="用户增长趋势",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="月度活跃用户数变化（单位：万）",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["1月", "2月", "3月", "4月", "5月", "6月",
                             "7月", "8月", "9月", "10月", "11月", "12月"]
    chart_data.add_series("2024年", (45, 48, 52, 58, 65, 72, 78, 85, 92, 98, 105, 118))
    chart_data.add_series("2025年", (52, 58, 65, 75, 88, 98, 110, 125, 138, 152, 165, 180))

    x, y, cx, cy = Inches(0.8), Inches(1.5), Inches(11.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = False
    plot.series[0].smooth = True
    plot.series[1].smooth = True

    # Reduce Y-axis tick density
    val_axis = chart.value_axis
    val_axis.major_unit = 40.0
    val_axis.minimum_scale = 0.0
    val_axis.maximum_scale = 200.0

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(11)

    return slide


def create_pie_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """饼图 - 市场份额分布"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="市场份额分布",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="2025年度各品牌市场占有率",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["品牌A", "品牌B", "品牌C", "品牌D", "其他"]
    chart_data.add_series("市场份额", (35, 28, 18, 12, 7))

    x, y, cx, cy = Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(11)

    add_text(
        slide, text="关键洞察：",
        left=6.5, top=2.0, width=6.0, height=0.4,
        font_size=16, bold=True, color="text", palette=palette,
    )

    insights = [
        "品牌A以35%份额领先市场",
        "品牌B保持稳健增长态势",
        "品牌C呈现上升趋势",
        "小品牌整合空间大",
    ]
    for i, insight in enumerate(insights):
        add_rect(
            slide,
            left=6.5, top=2.5 + i * 0.6, width=0.15, height=0.15,
            fill_color="primary", palette=palette,
        )
        add_text(
            slide, text=insight,
            left=6.8, top=2.45 + i * 0.6, width=6.0, height=0.5,
            font_size=13, bold=False, color="text", palette=palette,
        )

    return slide


def create_table_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """表格 - 产品功能对比"""
    from pptx.enum.text import PP_ALIGN

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="产品功能对比",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )

    rows, cols = 7, 4
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(5.5)
    ).table

    headers = ["功能特性", "基础版", "专业版", "企业版"]
    data = [
        ["用户数量限制", "5人", "50人", "不限"],
        ["存储空间", "10GB", "100GB", "1TB"],
        ["API调用", "1000次/日", "10000次/日", "不限"],
        ["数据分析", "基础报表", "高级分析", "全功能"],
        ["技术支持", "邮件支持", "7x24支持", "专属客服"],
        ["自定义域名", "不支持", "支持", "支持"],
    ]

    col_widths = [3.0, 3.0, 3.0, 3.333]
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    colors = BASE_PALETTES.get(palette, BASE_PALETTES["ocean_soft"])

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = rgb(colors["background"])
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(colors["primary"])

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = rgb(colors["text"])
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["light_bg"])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["background"])

    for i in range(rows):
        table.rows[i].height = Inches(0.75)

    return slide


def create_handdrawn_style_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """手绘风格 - AI产品开发流程数据"""
    import random

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="AI产品开发闭环",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="从需求洞察到持续迭代的完整数据驱动流程",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    def draw_jittered_box(bx, by, bw, bh, color):
        for i in range(15):
            t = i / 14
            x1 = bx + bw * t + random.uniform(-0.015, 0.015)
            y1 = by + random.uniform(-0.015, 0.015)
            x2 = bx + bw * (t + 1/14) + random.uniform(-0.015, 0.015)
            y2 = by + random.uniform(-0.015, 0.015)
            add_line(slide, x1, y1, x2, y2, color=color, width=1.5, palette=palette)
        for i in range(15):
            t = i / 14
            x1 = bx + bw * t + random.uniform(-0.015, 0.015)
            y1 = by + bh + random.uniform(-0.015, 0.015)
            x2 = bx + bw * (t + 1/14) + random.uniform(-0.015, 0.015)
            y2 = by + bh + random.uniform(-0.015, 0.015)
            add_line(slide, x1, y1, x2, y2, color=color, width=1.5, palette=palette)
        for i in range(15):
            t = i / 14
            x1 = bx + random.uniform(-0.015, 0.015)
            y1 = by + bh * t + random.uniform(-0.015, 0.015)
            x2 = bx + random.uniform(-0.015, 0.015)
            y2 = by + bh * (t + 1/14) + random.uniform(-0.015, 0.015)
            add_line(slide, x1, y1, x2, y2, color=color, width=1.5, palette=palette)
        for i in range(15):
            t = i / 14
            x1 = bx + bw + random.uniform(-0.015, 0.015)
            y1 = by + bh * t + random.uniform(-0.015, 0.015)
            x2 = bx + bw + random.uniform(-0.015, 0.015)
            y2 = by + bh * (t + 1/14) + random.uniform(-0.015, 0.015)
            add_line(slide, x1, y1, x2, y2, color=color, width=1.5, palette=palette)

    box_data = [
        {"x": 0.5, "title": "01 需求洞察", "color": "primary", "stats": [
            {"v": "150+", "l": "用户访谈"},
            {"v": "85%", "l": "需求命中率"},
            {"v": "3天", "l": "调研周期"},
        ]},
        {"x": 4.6, "title": "02 开发迭代", "color": "secondary", "stats": [
            {"v": "2周", "l": "Sprint周期"},
            {"v": "99.9%", "l": "代码通过率"},
            {"v": "50+", "l": "自动化用例"},
        ]},
        {"x": 8.7, "title": "03 效果验证", "color": "accent", "stats": [
            {"v": "4.8/5", "l": "用户满意度"},
            {"v": "300%", "l": "效率提升"},
            {"v": "60%", "l": "成本节省"},
        ]},
    ]

    for bd in box_data:
        draw_jittered_box(bd["x"], 1.55, 3.7, 3.5, bd["color"])
        add_text(
            slide, text=bd["title"],
            left=bd["x"] + 0.2, top=1.7, width=3.3, height=0.4,
            font_size=15, bold=True, color=bd["color"], alignment="left", palette=palette,
        )
        for i, stat in enumerate(bd["stats"]):
            sy = 2.2 + i * 0.95
            add_text(
                slide, text=stat["v"],
                left=bd["x"] + 0.2, top=sy, width=1.5, height=0.45,
                font_size=28, bold=True, color="text", alignment="left", palette=palette,
            )
            add_text(
                slide, text=stat["l"],
                left=bd["x"] + 0.2, top=sy + 0.45, width=3.3, height=0.3,
                font_size=11, bold=False, color="text_muted", alignment="left", palette=palette,
            )

    random.seed(100)
    for i in range(2):
        ax = box_data[i]["x"] + 3.7 + 0.1
        ay = 3.3
        bx = box_data[i + 1]["x"] - 0.1
        for seg in range(20):
            t = seg / 19
            x1 = ax + (bx - ax) * t
            y1 = ay + random.uniform(-0.015, 0.015) if seg > 0 else ay
            x2 = ax + (bx - ax) * (t + 1/20)
            y2 = ay + random.uniform(-0.015, 0.015)
            add_line(slide, x1, y1, x2, y2, color="divider", width=2.0, palette=palette)
        add_text(
            slide, text=">",
            left=(ax + bx) / 2 - 0.15, top=ay - 0.2, width=0.3, height=0.4,
            font_size=18, bold=True, color="divider", alignment="center", palette=palette,
        )

    add_text(
        slide, text="全流程周期：6-8周  |  持续迭代  |  数据驱动决策",
        left=0.5, top=6.55, width=12.833, height=0.3,
        font_size=12, bold=False, color="text_muted", alignment="center", palette=palette,
    )

    return slide


def create_icon_combo_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """图标组合 - 核心流程步骤"""
    import random
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="核心工作流程",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="AI驱动的端到端自动化流水线",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    steps = [
        {"icon": "01", "title": "需求采集", "desc": "多渠道收集用户需求，AI自动分类与优先级排序", "color": "primary"},
        {"icon": "02", "title": "数据分析", "desc": "结构化处理，提取关键指标与用户行为特征", "color": "secondary"},
        {"icon": "03", "title": "模型生成", "desc": "基于大模型自动生成多版本方案与评估报告", "color": "primary"},
        {"icon": "04", "title": "质量审核", "desc": "多模态QA自动检测，准确率提升至99%+", "color": "secondary"},
        {"icon": "05", "title": "优化迭代", "desc": "基于反馈数据持续优化，分钟级热更新", "color": "primary"},
        {"icon": "06", "title": "发布上线", "desc": "一键部署，自动灰度，全链路监控", "color": "secondary"},
    ]

    icon_size = 0.7
    for idx, step in enumerate(steps):
        row = idx // 3
        col = idx % 3
        x = 0.8 + col * 4.2
        y = 1.6 if row == 0 else 4.2

        add_ellipse(
            slide, left=x, top=y, width=icon_size, height=icon_size,
            fill_color=step["color"], palette=palette,
        )
        add_text(
            slide, text=step["icon"],
            left=x, top=y + 0.12, width=icon_size, height=0.45,
            font_size=18, bold=True, color="background", alignment="center", palette=palette,
        )
        add_text(
            slide, text=step["title"],
            left=x + icon_size + 0.2, top=y, width=2.8, height=0.4,
            font_size=16, bold=True, color="text", alignment="left", palette=palette,
        )
        add_text(
            slide, text=step["desc"],
            left=x + icon_size + 0.2, top=y + 0.38, width=2.8, height=0.5,
            font_size=11, bold=False, color="text_muted", alignment="left", palette=palette,
        )

        if col < 2:
            line_x1 = x + icon_size + 0.1
            line_y1 = y + icon_size / 2
            line_x2 = x + 4.2 - 0.1
            for seg in range(12):
                t = seg / 11
                x1 = line_x1 + (line_x2 - line_x1) * t
                y1 = line_y1 + random.uniform(-0.015, 0.015) if seg > 0 else line_y1
                x2 = line_x1 + (line_x2 - line_x1) * (t + 1/12)
                y2 = line_y1 + random.uniform(-0.015, 0.015)
                add_line(slide, x1, y1, x2, y2, color="divider", width=1.2, palette=palette)
            add_text(
                slide, text=">",
                left=(line_x1 + line_x2) / 2 - 0.12, top=line_y1 - 0.15,
                width=0.25, height=0.3,
                font_size=14, bold=True, color="divider", alignment="center", palette=palette,
            )

    add_rect(slide, left=0.5, top=6.35, width=12.333, height=0.5,
             fill_color="light_bg", palette=palette)
    footer_stats = [
        ("全流程自动化", "节省80%人力"),
        ("平均响应时间", "< 500ms"),
        ("服务可用性", "99.9%"),
        ("日均处理量", "1000万+"),
    ]
    for i, (label, val) in enumerate(footer_stats):
        fx = 0.7 + i * 3.1
        add_text(
            slide, text=label + "：",
            left=fx, top=6.4, width=1.4, height=0.35,
            font_size=10, bold=False, color="text_muted", alignment="left", palette=palette,
        )
        add_text(
            slide, text=val,
            left=fx + 1.3, top=6.4, width=1.6, height=0.35,
            font_size=11, bold=True, color="primary", alignment="left", palette=palette,
        )

    return slide


def create_smart_layout_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """智能布局 - 卡片网格+数据可视化+文字综合"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    add_text(
        slide, text="智能信息布局",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="卡片网格 + 数据可视化 + 文字说明的综合布局",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    # 顶部统计数据
    stats = [
        {"value": "98.5%", "label": "用户满意度", "color": "primary"},
        {"value": "2.3x", "label": "效率提升", "color": "secondary"},
        {"value": "50+", "label": "覆盖城市", "color": "accent"},
    ]
    for i, stat in enumerate(stats):
        x = 0.5 + i * 4.3
        add_rect(slide, left=x, top=1.4, width=4.0, height=1.3,
                 fill_color="light_bg", palette=palette)
        add_rect(slide, left=x, top=1.4, width=4.0, height=0.08,
                 fill_color=stat["color"], palette=palette)
        add_text(slide, text=stat["value"], left=x + 0.2, top=1.55, width=3.6, height=0.7,
                 font_size=32, bold=True, color=stat["color"], alignment="center", palette=palette)
        add_text(slide, text=stat["label"], left=x + 0.2, top=2.25, width=3.6, height=0.35,
                 font_size=12, color="text", alignment="center", palette=palette)

    # 左: 关键发现
    add_rect(slide, left=0.5, top=3.0, width=6.0, height=3.5,
             fill_color="background", palette=palette)
    add_text(slide, text="关键发现", left=0.7, top=3.15, width=5.6, height=0.4,
             font_size=16, bold=True, color="text", palette=palette)
    findings = [
        "移动端使用率首次超过PC端",
        "新功能用户留存率提升35%",
        "客户反馈响应时间缩短至2小时内",
    ]
    for i, finding in enumerate(findings):
        add_rect(slide, left=0.7, top=3.65 + i * 0.7, width=0.15, height=0.15,
                 fill_color="primary", palette=palette)
        add_text(slide, text=finding, left=1.0, top=3.6 + i * 0.7, width=5.3, height=0.6,
                 font_size=12, color="text", palette=palette)

    # 右: 趋势
    add_rect(slide, left=6.8, top=3.0, width=6.0, height=3.5,
             fill_color="light_bg", palette=palette)
    add_text(slide, text="趋势预览", left=7.0, top=3.15, width=5.6, height=0.4,
             font_size=16, bold=True, color="text", palette=palette)

    bar_data = [0.4, 0.55, 0.5, 0.7, 0.65, 0.85, 0.8]
    bar_width = 0.6
    gap = 0.15
    start_x = 7.2
    for i, height in enumerate(bar_data):
        x = start_x + i * (bar_width + gap)
        add_rect(slide, left=x, top=6.2 - height, width=bar_width, height=height,
                 fill_color="primary", palette=palette)

    add_text(slide, text="近7天活跃度", left=7.0, top=6.3, width=5.6, height=0.3,
             font_size=10, color="text_muted", palette=palette)

    return slide


# ============================================================================
# 主函数 - 生成完整示例
# ============================================================================

def generate_all_template_demo(palette: str = "ocean_soft", output_dir: str = "output"):
    """生成展示所有模板类型的完整PPT"""
    os.makedirs(output_dir, exist_ok=True)

    prs = new_presentation(palette=palette)
    slide_idx = 1

    def add_slide_note(title: str):
        nonlocal slide_idx
        print(f"[{slide_idx:02d}] {title}")
        slide_idx += 1

    # ================================================================
    # 第1页：封面
    # ================================================================
    add_slide_note("封面 - title_slide")
    generate_title_slide(
        prs, palette=palette,
        title="智能问答助手",
        subtitle="基于大模型的企业级AI交互平台",
        author="产品团队",
        date="2025年1月",
    )

    # ================================================================
    # 第2页：目录
    # ================================================================
    add_slide_note("目录 - agenda")
    generate_agenda(
        prs, palette=palette,
        kicker="目录",
        title="内容概览",
        items=[
            "01  产品概述与核心能力",
            "02  技术架构与创新",
            "03  应用场景与案例",
            "04  项目管理看板",
            "05  全国业务布局",
            "06  品牌价值主张",
            "07  效果数据与总结",
        ],
    )

# ============================================================================
# 主函数 - 生成完整示例
# ============================================================================

def generate_all_template_demo(palette: str = "ocean_soft", output_dir: str = "output"):
    """生成展示所有模板类型的完整PPT（纯结构预览，数据由用户填入）"""
    os.makedirs(output_dir, exist_ok=True)

    prs = new_presentation(palette=palette)
    slide_idx = 1

    def add_slide_note(title: str):
        nonlocal slide_idx
        print(f"[{slide_idx:02d}] {title}")
        slide_idx += 1

    # ================================================================
    # 第1页：封面
    # ================================================================
    add_slide_note("封面 - title_slide")
    generate_title_slide(prs, palette=palette)

    # ================================================================
    # 第2页：目录
    # ================================================================
    add_slide_note("目录 - agenda")
    generate_agenda(prs, palette=palette)

    # ================================================================
    # 第3页：大图介绍
    # ================================================================
    add_slide_note("大图介绍 - image_hero")
    generate_image_hero(prs, palette=palette)

    # ================================================================
    # 第4页：关键数字
    # ================================================================
    add_slide_note("关键数字 - stat_slide")
    generate_stat_slide(prs, palette=palette)

    # ================================================================
    # 第5页：章节分隔 - 产品概述
    # ================================================================
    add_slide_note("章节分隔 - section_divider")
    generate_section_divider(prs, palette=palette)

    # ================================================================
    # 第6页：内容页
    # ================================================================
    add_slide_note("内容页 - content_slide")
    generate_content_slide(prs, palette=palette)

    # ================================================================
    # 第7页：时间轴
    # ================================================================
    add_slide_note("时间轴 - timeline")
    generate_timeline(prs, palette=palette)

    # ================================================================
    # 第8页：流程图
    # ================================================================
    add_slide_note("流程图 - process_flow")
    generate_process_flow(prs, palette=palette)

    # ================================================================
    # 第9页：技术详解
    # ================================================================
    add_slide_note("技术详解 - deep_dive")
    generate_deep_dive(prs, palette=palette)

    # ================================================================
    # 第10页：章节分隔 - 核心功能
    # ================================================================
    add_slide_note("章节分隔 - section_divider")
    generate_section_divider(prs, palette=palette)

    # ================================================================
    # 第11页：卡片网格
    # ================================================================
    add_slide_note("卡片网格 - card_grid")
    generate_card_grid(prs, palette=palette)

    # ================================================================
    # 第12页：双栏对比
    # ================================================================
    add_slide_note("双栏对比 - two_column")
    generate_two_column(prs, palette=palette)

    # ================================================================
    # 第13页：图文混排
    # ================================================================
    add_slide_note("图文混排 - image_text")
    generate_image_text(prs, palette=palette)

    # ================================================================
    # 第14页：金句
    # ================================================================
    add_slide_note("金句 - quote_slide")
    generate_quote_slide(prs, palette=palette)

    # ================================================================
    # 第15页：章节分隔 - 行业应用
    # ================================================================
    add_slide_note("章节分隔 - section_divider")
    generate_section_divider(prs, palette=palette)

    # ================================================================
    # 第16页：三栏并列
    # ================================================================
    add_slide_note("三栏并列 - three_column")
    generate_three_column(prs, palette=palette)

    # ================================================================
    # 第17页：案例研究
    # ================================================================
    add_slide_note("案例研究 - case_study")
    generate_case_study(prs, palette=palette)

    # ================================================================
    # 第18页：KPI仪表盘
    # ================================================================
    add_slide_note("KPI仪表盘 - kpi_dashboard")
    generate_kpi_dashboard(prs, palette=palette)

    # ================================================================
    # 第19页：案例详情
    # ================================================================
    add_slide_note("案例详情 - example_detail")
    generate_example_detail(prs, palette=palette)

    # ================================================================
    # 第20页：看板
    # ================================================================
    add_slide_note("看板 - kanban")
    generate_kanban(prs, palette=palette)

    # ================================================================
    # 第21页：章节分隔 - 效果数据
    # ================================================================
    add_slide_note("章节分隔 - section_divider")
    generate_section_divider(prs, palette=palette)

    # ================================================================
    # 第22页：KPI仪表盘
    # ================================================================
    add_slide_note("KPI仪表盘 - kpi_dashboard")
    generate_kpi_dashboard(prs, palette=palette)

    # ================================================================
    # 第23页：柱状图
    # ================================================================
    add_slide_note("柱状图 - bar_chart")
    create_bar_chart_slide(prs, palette)

    # ================================================================
    # 第24页：折线图
    # ================================================================
    add_slide_note("折线图 - line_chart")
    create_line_chart_slide(prs, palette)

    # ================================================================
    # 第25页：饼图
    # ================================================================
    add_slide_note("饼图 - pie_chart")
    create_pie_chart_slide(prs, palette)

    # ================================================================
    # 第26页：表格
    # ================================================================
    add_slide_note("表格 - comparison_table")
    generate_comparison_table(prs, palette=palette)

    # ================================================================
    # 第27页：章节分隔 - 合作支持
    # ================================================================
    add_slide_note("章节分隔 - section_divider")
    generate_section_divider(prs, palette=palette)

    # ================================================================
    # 第28页：品牌聚焦
    # ================================================================
    add_slide_note("品牌聚焦 - brand_focus")
    generate_brand_focus(prs, palette=palette)

    # ================================================================
    # 第29页：图文混排
    # ================================================================
    add_slide_note("图文混排 - image_text")
    generate_image_text(prs, palette=palette)

    # ================================================================
    # 第30页：区域版图
    # ================================================================
    add_slide_note("区域版图 - region_map")
    generate_region_map(prs, palette=palette)

    # ================================================================
    # 第31页：总结
    # ================================================================
    add_slide_note("总结 - summary_slide")
    generate_summary_slide(prs, palette=palette)

    # 保存完整演示文稿
    output_path = os.path.join(output_dir, f"all_templates_{palette}.pptx")
    save_presentation(prs, output_path)
    print(f"\n完整演示文稿已保存: {output_path}")
    print(f"共 {slide_idx - 1} 页")

    # 单独保存图表类幻灯片（用于验证特殊元素）
    print(f"\n单独保存图表幻灯片...")
    chart_slides = [
        ("chart_01_bar", create_bar_chart_slide),
        ("chart_02_line", create_line_chart_slide),
        ("chart_03_pie", create_pie_chart_slide),
        ("chart_04_table", create_table_slide),
        ("chart_05_handdrawn", create_handdrawn_style_slide),
        ("chart_06_icon", create_icon_combo_slide),
        ("chart_07_smart", create_smart_layout_slide),
    ]
    for fname, func in chart_slides:
        prs_s = new_presentation(palette=palette)
        func(prs_s, palette)
        save_slide(prs_s.slides[0], os.path.join(output_dir, f"{fname}.pptx"))

    print(f"图表幻灯片已保存到 output/ 目录")

    return prs


if __name__ == "__main__":
    os.makedirs("output", exist_ok=True)

    print("=" * 60)
    print("生成所有模板类型示例")
    print("=" * 60)

    # ocean 配色
    generate_all_template_demo(palette="ocean_soft", output_dir="output")

    # warm 配色
    print()
    generate_all_template_demo(palette="warm_terracotta", output_dir="output")

    print()
    print("=" * 60)
    print("生成完成!")
    print("=" * 60)
