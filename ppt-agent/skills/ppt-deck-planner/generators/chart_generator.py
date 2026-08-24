#!/usr/bin/env python3
"""
数据可视化示例 - 展示图表和表格模板

这个脚本生成包含以下内容的示例PPT:
- 柱状图、折线图、饼图
- 表格
- 手绘风格页面

用法:
    python chart_generator.py

设计理念已整合到单页模板中:
- kanban_generator.py: 看板进度页
- brand_focus_generator.py: 品牌价值聚焦页
- region_map_generator.py: 区域版图页
"""
from __future__ import annotations
import os
import sys
import math
import random

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 复用项目的 base utilities
from generators.base import (
    PALETTES, rgb, new_presentation, set_slide_background,
    add_text, add_rect, add_round_rect, add_ellipse, add_line,
    add_source_line, save_presentation, save_slide,
)


# ============================================================================
# 方案一：数据可视化流 - 图表和表格
# ============================================================================

def create_bar_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建柱状图幻灯片 - 季度销售数据"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="季度销售业绩",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="各产品线季度收入对比（单位：万元）",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 一季度', 'Q2 二季度', 'Q3 三季度', 'Q4 四季度']
    chart_data.add_series('产品A线', (125, 148, 132, 189))
    chart_data.add_series('产品B线', (98, 115, 142, 167))
    chart_data.add_series('产品C线', (67, 78, 95, 112))

    x, y, cx, cy = Inches(0.8), Inches(1.5), Inches(11.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)

    return slide


def create_line_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建折线图幻灯片 - 趋势分析"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="用户增长趋势",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="月度活跃用户数变化（单位：万）",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ['1月', '2月', '3月', '4月', '5月', '6月',
                              '7月', '8月', '9月', '10月', '11月', '12月']
    chart_data.add_series('2024年', (45, 48, 52, 58, 65, 72, 78, 85, 92, 98, 105, 118))
    chart_data.add_series('2025年', (52, 58, 65, 75, 88, 98, 110, 125, 138, 152, 165, 180))

    x, y, cx, cy = Inches(0.8), Inches(1.5), Inches(11.0), Inches(5.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = False
    plot.series[0].smooth = True
    plot.series[1].smooth = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(11)

    return slide


def create_pie_chart_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建饼图幻灯片 - 市场份额"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="市场份额分布",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="2025年度各品牌市场占有率",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ['品牌A', '品牌B', '品牌C', '品牌D', '其他']
    chart_data.add_series('市场份额', (35, 28, 18, 12, 7))

    x, y, cx, cy = Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(11)

    add_text(
        slide, text="关键洞察：",
        left=6.5, top=2.0, width=6.0, height=0.4,
        font_size=16, bold=True, color="text", palette=palette,
    )

    insights = [
        "品牌A以35%份额领先市场",
        "品牌B保持稳健增长态势",
        "品牌C呈现上升趋势",
        "小品牌整合空间大"
    ]
    for i, insight in enumerate(insights):
        add_rect(
            slide,
            left=6.5, top=2.5 + i * 0.6, width=0.15, height=0.15,
            fill_color="primary", palette=palette,
        )
        add_text(
            slide, text=insight,
            left=6.8, top=2.45 + i * 0.6, width=6.0, height=0.5,
            font_size=13, bold=False, color="text", palette=palette,
        )

    return slide


def create_table_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建表格幻灯片 - 产品对比表"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="产品功能对比",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )

    rows, cols = 7, 4
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(5.5)
    ).table

    headers = ["功能特性", "基础版", "专业版", "企业版"]
    data = [
        ["用户数量限制", "5人", "50人", "不限"],
        ["存储空间", "10GB", "100GB", "1TB"],
        ["API调用", "1000次/日", "10000次/日", "不限"],
        ["数据分析", "基础报表", "高级分析", "全功能"],
        ["技术支持", "邮件支持", "7x24支持", "专属客服"],
        ["自定义域名", "不支持", "支持", "支持"],
    ]

    col_widths = [3.0, 3.0, 3.0, 3.333]
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = rgb(colors["background"])
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(colors["primary"])

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = rgb(colors["text"])
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["light_bg"])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["background"])

    for i in range(rows):
        table.rows[i].height = Inches(0.75)

    return slide


# ============================================================================
# 方案三：手绘风格流 - 手绘线条
# ============================================================================

def create_handdrawn_style_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建手绘风格幻灯片 - 流程思考"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="思考到行动",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="从问题分析到执行落地的完整闭环",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    # 手绘边框
    random.seed(456)
    box_left, box_top = 1.0, 1.5
    box_right, box_bottom = 12.333, 6.8

    for i in range(20):
        x1 = box_left + (box_right - box_left) * i / 20
        y1 = box_top + random.uniform(-0.03, 0.03)
        x2 = box_left + (box_right - box_left) * (i + 1) / 20
        y2 = box_top + random.uniform(-0.03, 0.03)
        add_line(slide, x1, y1, x2, y2, color="primary", width=1.5, palette=palette)

    for i in range(20):
        x1 = box_left + (box_right - box_left) * i / 20
        y1 = box_bottom + random.uniform(-0.03, 0.03)
        x2 = box_left + (box_right - box_left) * (i + 1) / 20
        y2 = box_bottom + random.uniform(-0.03, 0.03)
        add_line(slide, x1, y1, x2, y2, color="primary", width=1.5, palette=palette)

    for i in range(20):
        x1 = box_left + random.uniform(-0.03, 0.03)
        y1 = box_top + (box_bottom - box_top) * i / 20
        x2 = box_left + random.uniform(-0.03, 0.03)
        y2 = box_top + (box_bottom - box_top) * (i + 1) / 20
        add_line(slide, x1, y1, x2, y2, color="primary", width=1.5, palette=palette)

    for i in range(20):
        x1 = box_right + random.uniform(-0.03, 0.03)
        y1 = box_top + (box_bottom - box_top) * i / 20
        x2 = box_right + random.uniform(-0.03, 0.03)
        y2 = box_top + (box_bottom - box_top) * (i + 1) / 20
        add_line(slide, x1, y1, x2, y2, color="primary", width=1.5, palette=palette)

    # 手绘圆形
    random.seed(789)
    circle_x, circle_y = 2.5, 4.0
    circle_r = 1.2
    for angle in range(0, 360, 5):
        rad = math.radians(angle)
        jitter_r = circle_r + random.uniform(-0.05, 0.05)
        x1 = circle_x + jitter_r * math.cos(rad)
        y1 = circle_y + jitter_r * math.sin(rad)
        x2 = circle_x + jitter_r * math.cos(rad + math.radians(6))
        y2 = circle_y + jitter_r * math.sin(rad + math.radians(6))
        add_line(slide, x1, y1, x2, y2, color="secondary", width=1.5, palette=palette)

    add_text(
        slide, text="思考",
        left=2.0, top=3.6, width=1.0, height=0.5,
        font_size=14, bold=True, color="text", palette=palette,
    )

    # 手绘箭头
    random.seed(321)
    arrow_start = (4.5, 4.0)
    arrow_end = (6.5, 3.5)
    for i in range(15):
        t = i / 14
        x1 = arrow_start[0] + (arrow_end[0] - arrow_start[0]) * t + random.uniform(-0.04, 0.04)
        y1 = arrow_start[1] + (arrow_end[1] - arrow_start[1]) * t + random.uniform(-0.04, 0.04)
        x2 = arrow_start[0] + (arrow_end[0] - arrow_start[0]) * (t + 1/14) + random.uniform(-0.04, 0.04)
        y2 = arrow_start[1] + (arrow_end[1] - arrow_start[1]) * (t + 1/14) + random.uniform(-0.04, 0.04)
        add_line(slide, x1, y1, x2, y2, color="accent", width=2.0, palette=palette)

    add_text(
        slide, text=">",
        left=6.2, top=3.2, width=0.5, height=0.5,
        font_size=24, bold=True, color="accent", palette=palette,
    )

    # 手绘矩形
    random.seed(654)
    rect_x, rect_y = 7.5, 2.8
    rect_w, rect_h = 2.5, 2.0
    corners = [
        (rect_x, rect_y),
        (rect_x + rect_w, rect_y),
        (rect_x + rect_w, rect_y + rect_h),
        (rect_x, rect_y + rect_h),
        (rect_x, rect_y),
    ]
    for i in range(4):
        for j in range(10):
            t = j / 10
            x1 = corners[i][0] + (corners[i+1][0] - corners[i][0]) * t + random.uniform(-0.03, 0.03)
            y1 = corners[i][1] + (corners[i+1][1] - corners[i][1]) * t + random.uniform(-0.03, 0.03)
            x2 = corners[i][0] + (corners[i+1][0] - corners[i][0]) * (t + 1/10) + random.uniform(-0.03, 0.03)
            y2 = corners[i][1] + (corners[i+1][1] - corners[i][1]) * (t + 1/10) + random.uniform(-0.03, 0.03)
            add_line(slide, x1, y1, x2, y2, color="primary", width=1.5, palette=palette)

    add_text(
        slide, text="行动",
        left=8.0, top=3.6, width=1.5, height=0.5,
        font_size=14, bold=True, color="text", palette=palette,
    )

    add_text(
        slide, text="注：python-pptx 原生不支持 sketch_format，此效果通过随机抖动线段坐标实现",
        left=1.0, top=6.9, width=11.0, height=0.3,
        font_size=10, bold=False, color="text_muted", palette=palette,
    )

    return slide


# ============================================================================
# 方案四：视觉分层流 - 图标组合和信息布局
# ============================================================================

def create_icon_combo_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建图标组合幻灯片"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="图标组合设计",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="用圆形、矩形等基础形状组合成复杂视觉元素",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    # 图标1: 用户
    icon1_x, icon1_y = 1.5, 2.5
    add_ellipse(slide, icon1_x + 0.4, icon1_y, 0.5, 0.5, fill_color="primary", palette=palette)
    add_round_rect(slide, icon1_x + 0.15, icon1_y + 0.6, 1.0, 0.8,
                   fill_color="primary", palette=palette)
    add_text(slide, text="用户", left=icon1_x, top=icon1_y + 1.5, width=1.3, height=0.4,
             font_size=12, color="text", alignment="center", palette=palette)

    # 图标2: 数据
    icon2_x, icon2_y = 4.0, 2.5
    add_rect(slide, icon2_x, icon2_y, 1.3, 1.3, fill_color="secondary", palette=palette)
    add_rect(slide, icon2_x + 0.15, icon2_y + 0.5, 0.25, 0.6, fill_color="background", palette=palette)
    add_rect(slide, icon2_x + 0.5, icon2_y + 0.25, 0.25, 0.85, fill_color="background", palette=palette)
    add_rect(slide, icon2_x + 0.85, icon2_y + 0.65, 0.25, 0.45, fill_color="background", palette=palette)
    add_text(slide, text="数据", left=icon2_x, top=icon2_y + 1.5, width=1.3, height=0.4,
             font_size=12, color="text", alignment="center", palette=palette)

    # 图标3: 设置
    icon3_x, icon3_y = 6.5, 2.5
    add_ellipse(slide, icon3_x + 0.15, icon3_y + 0.15, 1.0, 1.0, fill_color="accent", palette=palette)
    for angle in range(0, 360, 45):
        rad = math.radians(angle)
        tx = icon3_x + 0.65 + 0.55 * math.cos(rad) - 0.08
        ty = icon3_y + 0.65 + 0.55 * math.sin(rad) - 0.04
        add_rect(slide, tx, ty, 0.16, 0.25, fill_color="accent", palette=palette)
    add_ellipse(slide, icon3_x + 0.4, icon3_y + 0.4, 0.5, 0.5, fill_color="background", palette=palette)
    add_text(slide, text="设置", left=icon3_x, top=icon3_y + 1.5, width=1.3, height=0.4,
             font_size=12, color="text", alignment="center", palette=palette)

    # 图标4: 消息
    icon4_x, icon4_y = 9.0, 2.5
    add_round_rect(slide, icon4_x, icon4_y + 0.15, 1.2, 0.9, fill_color="primary", palette=palette)
    add_rect(slide, icon4_x + 0.1, icon4_y + 0.95, 0.3, 0.25,
             fill_color="primary", palette=palette)
    add_rect(slide, icon4_x + 0.15, icon4_y + 0.35, 0.9, 0.12, fill_color="background", palette=palette)
    add_rect(slide, icon4_x + 0.15, icon4_y + 0.55, 0.7, 0.12, fill_color="background", palette=palette)
    add_rect(slide, icon4_x + 0.15, icon4_y + 0.75, 0.5, 0.12, fill_color="background", palette=palette)
    add_text(slide, text="消息", left=icon4_x, top=icon4_y + 1.5, width=1.3, height=0.4,
             font_size=12, color="text", alignment="center", palette=palette)

    # 底部说明
    add_rect(slide, left=0.5, top=5.5, width=12.333, height=1.5,
             fill_color="light_bg", palette=palette)
    add_text(
        slide, text="设计要点：",
        left=0.8, top=5.65, width=11.0, height=0.35,
        font_size=14, bold=True, color="text", palette=palette,
    )
    tips = [
        "使用相同的设计语言（圆角、颜色）保持一致性",
        "形状的堆叠顺序决定视觉层次",
        "适当留白让图标更清晰易读",
    ]
    for i, tip in enumerate(tips):
        add_rect(slide, left=0.8, top=6.05 + i * 0.3, width=0.1, height=0.1,
                 fill_color="secondary", palette=palette)
        add_text(slide, text=tip, left=1.0, top=6.0 + i * 0.3, width=11.0, height=0.3,
                 font_size=11, color="text_muted", palette=palette)

    return slide


def create_smart_layout_slide(prs, palette: str = "ocean_soft") -> Presentation:
    """创建智能布局幻灯片"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, palette)

    colors = PALETTES.get(palette, PALETTES["ocean_soft"])

    add_text(
        slide, text="智能信息布局",
        left=0.5, top=0.3, width=12.0, height=0.6,
        font_size=32, bold=True, color="text", palette=palette,
    )
    add_text(
        slide, text="卡片网格 + 数据可视化 + 文字说明的综合布局",
        left=0.5, top=0.85, width=12.0, height=0.35,
        font_size=14, bold=False, color="text_muted", palette=palette,
    )

    # 顶部统计数据
    stats = [
        {"value": "98.5%", "label": "用户满意度", "color": "primary"},
        {"value": "2.3x", "label": "效率提升", "color": "secondary"},
        {"value": "50+", "label": "覆盖城市", "color": "accent"},
    ]
    for i, stat in enumerate(stats):
        x = 0.5 + i * 4.3
        add_rect(slide, left=x, top=1.4, width=4.0, height=1.3,
                 fill_color="light_bg", palette=palette)
        add_rect(slide, left=x, top=1.4, width=4.0, height=0.08,
                 fill_color=stat["color"], palette=palette)
        add_text(slide, text=stat["value"], left=x + 0.2, top=1.55, width=3.6, height=0.7,
                 font_size=32, bold=True, color=stat["color"], alignment="center", palette=palette)
        add_text(slide, text=stat["label"], left=x + 0.2, top=2.25, width=3.6, height=0.35,
                 font_size=12, color="text", alignment="center", palette=palette)

    # 左: 关键发现
    add_rect(slide, left=0.5, top=3.0, width=6.0, height=3.5,
             fill_color="background", palette=palette)
    add_text(slide, text="关键发现", left=0.7, top=3.15, width=5.6, height=0.4,
             font_size=16, bold=True, color="text", palette=palette)
    findings = [
        "移动端使用率首次超过PC端",
        "新功能用户留存率提升35%",
        "客户反馈响应时间缩短至2小时内",
    ]
    for i, finding in enumerate(findings):
        add_rect(slide, left=0.7, top=3.65 + i * 0.7, width=0.15, height=0.15,
                 fill_color="primary", palette=palette)
        add_text(slide, text=finding, left=1.0, top=3.6 + i * 0.7, width=5.3, height=0.6,
                 font_size=12, color="text", palette=palette)

    # 右: 趋势
    add_rect(slide, left=6.8, top=3.0, width=6.0, height=3.5,
             fill_color="light_bg", palette=palette)
    add_text(slide, text="趋势预览", left=7.0, top=3.15, width=5.6, height=0.4,
             font_size=16, bold=True, color="text", palette=palette)

    bar_data = [0.4, 0.55, 0.5, 0.7, 0.65, 0.85, 0.8]
    bar_width = 0.6
    gap = 0.15
    start_x = 7.2
    for i, height in enumerate(bar_data):
        x = start_x + i * (bar_width + gap)
        add_rect(slide, left=x, top=6.2 - height, width=bar_width, height=height,
                 fill_color="primary", palette=palette)

    add_text(slide, text="近7天活跃度", left=7.0, top=6.3, width=5.6, height=0.3,
             font_size=10, color="text_muted", palette=palette)

    return slide


# ============================================================================
# 主函数 - 生成示例
# ============================================================================

def generate_all_examples(output_dir: str = "output"):
    """生成所有示例幻灯片"""
    os.makedirs(output_dir, exist_ok=True)

    palette = "ocean_soft"
    prs = new_presentation(palette=palette)

    print("=" * 60)
    print("数据可视化示例")
    print("=" * 60)

    create_bar_chart_slide(prs, palette)
    print("[1/7] 柱状图幻灯片 - 已创建")

    create_line_chart_slide(prs, palette)
    print("[2/7] 折线图幻灯片 - 已创建")

    create_pie_chart_slide(prs, palette)
    print("[3/7] 饼图幻灯片 - 已创建")

    create_table_slide(prs, palette)
    print("[4/7] 表格幻灯片 - 已创建")

    create_handdrawn_style_slide(prs, palette)
    print("[5/7] 手绘风格幻灯片 - 已创建")

    create_icon_combo_slide(prs, palette)
    print("[6/7] 图标组合幻灯片 - 已创建")

    create_smart_layout_slide(prs, palette)
    print("[7/7] 智能布局幻灯片 - 已创建")

    # 保存
    output_path = os.path.join(output_dir, "data_visualization_examples.pptx")
    save_presentation(prs, output_path)
    print(f"\n完整演示文稿已保存: {output_path}")

    # 单独保存
    slides = [
        ("01_bar_chart.pptx", create_bar_chart_slide),
        ("02_line_chart.pptx", create_line_chart_slide),
        ("03_pie_chart.pptx", create_pie_chart_slide),
        ("04_table.pptx", create_table_slide),
        ("05_handdrawn.pptx", create_handdrawn_style_slide),
        ("06_icon_combo.pptx", create_icon_combo_slide),
        ("07_smart_layout.pptx", create_smart_layout_slide),
    ]

    for fname, func in slides:
        prs_s = new_presentation(palette=palette)
        func(prs_s, palette)
        save_slide(prs_s.slides[0], os.path.join(output_dir, fname))
        print(f"单独幻灯片已保存: output/{fname}")

    print("\n" + "=" * 60)
    print("生成完成!")
    print("=" * 60)
    print("\n新设计理念已整合到单页模板中:")
    print("  - kanban_generator.py: 看板进度页")
    print("  - brand_focus_generator.py: 品牌价值聚焦页")
    print("  - region_map_generator.py: 区域版图页")
    print("=" * 60)

    return prs


if __name__ == "__main__":
    generate_all_examples()
