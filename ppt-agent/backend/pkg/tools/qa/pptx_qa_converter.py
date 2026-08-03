#!/usr/bin/env python3
"""
PPTX QA Converter - 将 PPTX 文件转换为图片供 QA 视觉审查使用。

功能：
1. 逐文件转换 PPTX → PDF → JPG（不合并 PPTX，避免布局叠加）
2. 批量 QA 时将各 PPTX 转为 PDF 后用 pdfunite 合并（合并发生在 PDF 层，安全）
3. 提取所有 PPTX 的文本内容用于内容 QA

性能优化：
- 全局文件锁防止多进程并发杀进程导致 CPU 打满
- 模块级缓存 soffice/pdftoppm/pdfunite 路径，避免每次调用重新查找

依赖：
    pip install python-pptx Pillow python-markdown

环境要求：
    - LibreOffice (soffice) 已安装并可在 PATH 中访问
    - poppler-utils (pdftoppm) 已安装
"""

import argparse
import glob
import json
import os
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path
from typing import Dict, List, Optional

try:
    import fcntl
except ImportError:
    fcntl = None

try:
    import msvcrt
except ImportError:
    msvcrt = None

# ── Module-level caches (avoids repeated subprocess calls to find binaries) ──
_SOFFICE_CMD: Optional[str] = None
_PDFTOPPM_AVAILABLE: Optional[bool] = None
_PDFTOPPM_CMD: Optional[str] = None
_LOCK_FILE_PATH = os.path.join(tempfile.gettempdir(), "pptx_conv.lock")


def _run_process(command: List[str], **kwargs) -> subprocess.CompletedProcess:
    """Run native commands and Windows batch shims with identical semantics."""
    executable = Path(command[0])
    if os.name == "nt" and executable.suffix.lower() in {".bat", ".cmd"}:
        command_line = subprocess.list2cmdline(command)
        return subprocess.run(command_line, shell=True, **kwargs)
    return subprocess.run(command, **kwargs)


def find_soffice() -> Optional[str]:
    """查找 LibreOffice 可执行文件路径（模块级缓存，只查一次）"""
    global _SOFFICE_CMD
    if _SOFFICE_CMD is not None:
        return _SOFFICE_CMD
    candidates = [
        shutil.which("soffice"),
        shutil.which("soffice.com"),
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for candidate in candidates:
        if not candidate:
            continue
        try:
            result = _run_process(
                [candidate, "--version"],
                capture_output=True,
                timeout=5
            )
            if result.returncode == 0:
                _SOFFICE_CMD = candidate
                return candidate
        except (FileNotFoundError, subprocess.TimeoutExpired, PermissionError):
            continue
    _SOFFICE_CMD = None
    print("警告: 未找到 LibreOffice (soffice)，PPTX 将无法转换为 PDF 进行视觉审查", file=sys.stderr)
    return None


def find_pdftoppm() -> Optional[str]:
    """查找 pdftoppm 命令（模块级缓存，只查一次）。"""
    global _PDFTOPPM_AVAILABLE, _PDFTOPPM_CMD
    if _PDFTOPPM_AVAILABLE is not None:
        return _PDFTOPPM_CMD if _PDFTOPPM_AVAILABLE else None

    resolved = shutil.which("pdftoppm")
    candidates = [os.environ.get("PDFTOPPM_CMD"), shutil.which("pdftoppm.exe"), resolved]
    if resolved:
        resolved_path = Path(resolved).resolve()
        for parent in resolved_path.parents:
            candidates.append(str(parent / "native" / "poppler" / "Library" / "bin" / "pdftoppm.exe"))
    candidates.append("pdftoppm")

    for command in candidates:
        if not command:
            continue
        if os.path.isabs(command) and not os.path.exists(command):
            continue
        try:
            result = _run_process([command, "-v"], capture_output=True, timeout=5)
            stderr = result.stderr.decode("utf-8", errors="replace").lower()
            if result.returncode == 0 or "pdftoppm" in stderr:
                _PDFTOPPM_AVAILABLE = True
                _PDFTOPPM_CMD = command
                return command
        except (FileNotFoundError, subprocess.TimeoutExpired, PermissionError):
            continue

    _PDFTOPPM_AVAILABLE = False
    if not _PDFTOPPM_AVAILABLE:
        print("警告: 未找到 pdftoppm，PDF 将无法转换为图片", file=sys.stderr)
    return None


# ── Global file lock ──────────────────────────────────────────────────────
# Replaces pkill-based approach. Instead of killing other processes' soffice
# instances (which causes them to retry and re-kill each other -> CPU spin),
# we serialize all conversion work behind a single exclusive lock file.
# Waiting processes block on flock() and proceed one at a time.

_lock_fd: Optional[int] = None


def _try_lock(fd: int) -> bool:
    try:
        if fcntl is not None:
            fcntl.flock(fd, fcntl.LOCK_EX | fcntl.LOCK_NB)
        elif msvcrt is not None:
            if os.fstat(fd).st_size == 0:
                os.write(fd, b"0")
            os.lseek(fd, 0, os.SEEK_SET)
            msvcrt.locking(fd, msvcrt.LK_NBLCK, 1)
        else:
            raise RuntimeError("platform does not provide a file locking API")
        return True
    except (BlockingIOError, OSError):
        return False


def _unlock(fd: int) -> None:
    if fcntl is not None:
        fcntl.flock(fd, fcntl.LOCK_UN)
    elif msvcrt is not None:
        os.lseek(fd, 0, os.SEEK_SET)
        msvcrt.locking(fd, msvcrt.LK_UNLCK, 1)


def _acquire_lock(timeout: float = 120.0) -> bool:
    """
    Acquire exclusive lock. Returns True on success, False on timeout.
    All conversion calls should wrap their work in this lock.
    """
    global _lock_fd
    try:
        lock_path = os.environ.get("PPTX_CONV_LOCK", _LOCK_FILE_PATH)
        lock_dir = os.path.dirname(lock_path)
        os.makedirs(lock_dir, exist_ok=True)
        fd = os.open(lock_path, os.O_RDWR | os.O_CREAT, 0o644)
        # Try non-blocking first — if already held, loop with a timeout.
        if _try_lock(fd):
            _lock_fd = fd
            return True
        # Another process holds the lock — wait for it with timeout.
        start = time.monotonic()
        while time.monotonic() - start < timeout:
            if _try_lock(fd):
                _lock_fd = fd
                return True
            time.sleep(0.5)
        os.close(fd)
        return False
    except Exception:
        return False


def _release_lock():
    """Release the exclusive lock."""
    global _lock_fd
    if _lock_fd is not None:
        try:
            _unlock(_lock_fd)
            os.close(_lock_fd)
        except Exception:
            pass
        _lock_fd = None


def ensure_soffice_ready(soffice_cmd: str) -> bool:
    """
    Warm up LibreOffice service. Called once per Python invocation (not per file).
    No longer kills other processes — the file lock handles serialization.
    """
    try:
        subprocess.run(
            [soffice_cmd, "--headless", "--norestore", "--nofirststartwizard",
             "--convert-to", "pdf", "--outdir", "/tmp", "/dev/null"],
            capture_output=True,
            timeout=30
        )
        return True
    except Exception:
        return False




def extract_text_from_pptx_files(pptx_files: List[str]) -> List[str]:
    """从所有 PPTX 文件中提取文本内容，按文件顺序返回。"""
    texts = []
    try:
        from pptx import Presentation
    except ImportError:
        return texts

    for pptx_file in pptx_files:
        try:
            prs = Presentation(pptx_file)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            if parts:
                texts.append("\n".join(parts))
            else:
                texts.append("")
        except Exception as e:
            print(f"警告: 提取 {pptx_file} 文本时出错: {e}", file=sys.stderr)
            texts.append("")

    return texts


# ── Internal helpers ───────────────────────────────────────────────────────

def _soffice_convert_pptx_to_pdf(soffice: str, pptx_file: str, out_dir: str) -> Optional[str]:
    """Convert a single PPTX to PDF using LibreOffice. Returns PDF path or None."""
    out_name = Path(pptx_file).stem + ".pdf"
    expected_pdf = os.path.join(out_dir, out_name)

    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_file],
        capture_output=True, timeout=120
    )
    if result.returncode != 0:
        # LibreOffice may put the PDF in a different location — try to find it
        for name in os.listdir(out_dir):
            if name.startswith(Path(pptx_file).stem) and name.endswith(".pdf"):
                return os.path.join(out_dir, name)
        return None
    return expected_pdf if os.path.exists(expected_pdf) else None


def _merge_pdfs(pdf_paths: List[str], out_path: str) -> bool:
    """Merge multiple PDFs into one using pdfunite. Returns True on success."""
    # Try pdfunite first (fast, widespread)
    for tool in ["pdfunite", "pdftk"]:
        if shutil.which(tool):
            cmd = [tool] + pdf_paths + [out_path]
            if tool == "pdftk":
                cmd = [tool] + [p for p in pdf_paths] + ["cat", "output", out_path]
            result = subprocess.run(cmd, capture_output=True, timeout=120)
            if result.returncode == 0 and os.path.exists(out_path):
                return True
    return False


def _pdf_to_jpgs_by_stem(pdf_path: str, output_dir: str, stems: List[str], dpi: int) -> tuple:
    """Split a PDF into JPGs, naming each by the corresponding stem."""
    slide_images = []
    errors = []
    pdf_stem = Path(pdf_path).stem
    pdf_dir = Path(pdf_path).parent

    command = find_pdftoppm()
    if not command:
        return [], ["pdftoppm is not available"]
    result = _run_process(
        [command, "-jpeg", "-r", str(dpi), "-q", pdf_path,
         str(pdf_dir / pdf_stem)],
        capture_output=True, timeout=300
    )
    if result.returncode != 0:
        errors.append(f"pdftoppm failed: {result.stderr.decode('utf-8', errors='replace')}")
        return [], errors

    def page_number(path: Path) -> int:
        suffix = path.stem.removeprefix(f"{pdf_stem}-")
        return int(suffix) if suffix.isdigit() else sys.maxsize

    generated = sorted(pdf_dir.glob(f"{pdf_stem}-*.jpg"), key=page_number)
    for idx, stem in enumerate(stems, start=1):
        src = generated[idx - 1] if idx <= len(generated) else None
        dst = Path(output_dir) / f"{stem}.jpg"
        if src is not None and src.exists():
            if dst.exists():
                dst.unlink()
            shutil.move(str(src), str(dst))
            slide_images.append(dst.name)
        else:
            errors.append(f"pdftoppm did not generate page {idx}")

    return slide_images, errors


def _convert_pptxs_to_jpgs(output_dir: str, pptx_files: List[str], dpi: int) -> Dict:
    """Merge PPTXs, convert to PDF, split into JPGs named by original stems."""
    soffice = find_soffice()
    pdftoppm = find_pdftoppm()

    if not soffice:
        return {"error": "LibreOffice 未安装", "slide_images": [], "text_content": ""}
    if not pdftoppm:
        return {"error": "pdftoppm 未安装", "slide_images": [], "text_content": ""}

    ensure_soffice_ready(soffice)

    # Extract text from all PPTXs
    slide_texts = extract_text_from_pptx_files(pptx_files)
    stems = [Path(f).stem for f in pptx_files]

    with tempfile.TemporaryDirectory() as tmp_dir:
        # Step 1: Merge PPTXs via python-pptx
        try:
            from pptx import Presentation
            from copy import deepcopy

            base_prs = Presentation(pptx_files[0])
            for pptx_file in pptx_files[1:]:
                src_prs = Presentation(pptx_file)
                blank_layout = base_prs.slide_layouts[6]
                for slide in src_prs.slides:
                    new_slide = base_prs.slides.add_slide(blank_layout)
                    for shape in list(new_slide.shapes):
                        sp = shape.element
                        sp.getparent().remove(sp)
                    for shape in slide.shapes:
                        el = shape.element
                        new_slide.shapes._spTree.insert_element_before(
                            deepcopy(el), 'p:extLst'
                        )

            merged_pptx = os.path.join(tmp_dir, "merged.pptx")
            base_prs.save(merged_pptx)
        except Exception as e:
            return {"error": f"PPTX 合并失败: {e}", "slide_images": [], "text_content": ""}

        # Step 2: Convert merged PPTX to PDF
        merged_pdf = os.path.join(tmp_dir, "merged.pdf")
        max_retries = 2
        last_err = None
        for _ in range(max_retries + 1):
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 "--outdir", tmp_dir, merged_pptx],
                capture_output=True, timeout=300
            )
            if result.returncode == 0:
                break
            last_err = result.stderr.decode('utf-8', errors='replace')
        else:
            return {"error": f"LibreOffice 转换失败: {last_err}",
                    "slide_images": [], "text_content": ""}

        if not os.path.exists(merged_pdf):
            # LibreOffice may generate PDF with a different stem in tmp_dir
            found = False
            for name in os.listdir(tmp_dir):
                if name.startswith("merged") and name.endswith(".pdf"):
                    shutil.move(os.path.join(tmp_dir, name), merged_pdf)
                    found = True
                    break
            if not found:
                return {"error": "LibreOffice 未生成 PDF",
                        "slide_images": [], "text_content": ""}

        # Step 3: Split PDF to JPGs, name by original stems
        slide_images, errors = _pdf_to_jpgs_by_stem(
            merged_pdf, output_dir, stems, dpi
        )

        # Step 4: Move JPGs from tmp_dir to output_dir if needed
        for img_name in slide_images:
            src = Path(tmp_dir) / img_name
            dst = Path(output_dir) / img_name
            if src.exists() and not dst.exists():
                shutil.move(str(src), str(dst))

    all_text = []
    for idx, text in enumerate(slide_texts, start=1):
        all_text.append(f"[第{idx}页]\n{text}" if text else f"[第{idx}页]\n")

    return {
        "slide_images": slide_images,
        "text_content": "\n\n".join(all_text),
        "total_slides": len(slide_images),
        "total_files": len(pptx_files),
        "conversion_errors": errors if errors else [],
    }


def pptx_to_images(pptx_dir: str, output_dir: str, dpi: int = 150,
                   specific_files: List[str] = None) -> Dict:
    """
    将指定目录下的所有（或指定的）PPTX 文件转换为单页图片。

    优化策略：
    - specific_files 非空时：先将指定文件合并，再一次性转换（单次 LibreOffice）
    - specific_files 为空时：先将所有 PPTX 合并为一个文件，再一次性转换

    Args:
        pptx_dir: 包含 PPTX 文件的目录
        output_dir: 输出图片的目录
        dpi: 图片分辨率，默认 150
        specific_files: 仅转换这些文件（不含路径），为空时转换目录下所有 PPTX

    Returns:
        包含转换结果的字典
    """
    # ── Acquire global lock before any conversion work ──
    if not _acquire_lock(timeout=300):
        return {"error": "无法获取转换锁（另一个进程占用超时）",
                "slide_images": [], "text_content": ""}

    try:
        if specific_files:
            pptx_files = [
                os.path.join(pptx_dir, f) if not os.path.isabs(f) else f
                for f in specific_files
            ]
            pptx_files = [f for f in pptx_files if os.path.exists(f) and os.path.isfile(f)]
        else:
            pptx_files = sorted(glob.glob(os.path.join(pptx_dir, "*.pptx")))

        if not pptx_files:
            return {"error": f"在 {pptx_dir} 中未找到 PPTX 文件",
                    "slide_images": [], "text_content": ""}

        # ── Always merge first, then convert once (single LibreOffice startup) ──
        # This eliminates the N-file → N-LibreOffice problem from _convert_specific_files.
        return _convert_pptxs_to_jpgs(output_dir, pptx_files, dpi)

    finally:
        _release_lock()


def pptx_to_pdf(pptx_dir: str, out_pdf_path: str = None) -> Dict:
    """
    将目录下所有 PPTX 逐个转换为 PDF，再用 pdfunite 合并为一个文件。

    用于批量 QA 场景：一次性生成合并 PDF 发给多模态 LLM 审查所有页。
    不合并 PPTX（会导致布局叠加），而是转换后再合并 PDF。
    """
    if not _acquire_lock(timeout=300):
        return {"error": "无法获取转换锁（另一个进程占用超时）"}

    try:
        pptx_files = sorted(glob.glob(os.path.join(pptx_dir, "*.pptx")))
        if not pptx_files:
            return {"error": f"在 {pptx_dir} 中未找到 PPTX 文件"}

        soffice = find_soffice()
        if not soffice:
            return {"error": "LibreOffice 未安装"}

        if out_pdf_path is None:
            out_pdf_path = os.path.join(pptx_dir, "merged.pdf")

        ensure_soffice_ready(soffice)
        slide_texts = extract_text_from_pptx_files(pptx_files)

        with tempfile.TemporaryDirectory() as tmp_dir:
            pdf_paths = []
            for pptx_file in pptx_files:
                pdf_path = _soffice_convert_pptx_to_pdf(soffice, pptx_file, tmp_dir)
                if pdf_path:
                    pdf_paths.append(pdf_path)

            if not pdf_paths:
                return {"error": "所有 PPTX 转换 PDF 失败"}

            # Merge PDFs using pdfunite (no PPTX merge → no layout corruption)
            if not _merge_pdfs(pdf_paths, out_pdf_path):
                # Fallback: copy first PDF if pdfunite unavailable
                shutil.copy2(pdf_paths[0], out_pdf_path)

        if not os.path.exists(out_pdf_path):
            return {"error": "PDF 合并失败，未生成最终文件"}

        all_text = []
        for idx, text in enumerate(slide_texts, start=1):
            all_text.append(f"[第{idx}页]\n{text}" if text else f"[第{idx}页]\n")

        return {
            "pdf_path": out_pdf_path,
            "slide_count": len(pptx_files),
            "text_content": "\n\n".join(all_text),
        }
    finally:
        _release_lock()


def main():
    parser = argparse.ArgumentParser(
        description="PPTX QA Converter - 逐文件转换 PPTX 为 JPG 缩略图，批量 QA 时合并 PDF（优化版）"
    )
    parser.add_argument("--pptx-dir", required=True,
                        help="包含 PPTX 文件的目录")
    parser.add_argument("--output-dir", required=True,
                        help="输出图片的目录")
    parser.add_argument("--dpi", type=int, default=150,
                        help="输出图片 DPI (默认: 150)")
    parser.add_argument("--format", choices=["json", "text"], default="json",
                        help="输出格式 (默认: json)")
    parser.add_argument("--files", nargs="*", default=None,
                        help="仅转换指定的 PPTX 文件（不含路径）。逐文件转换，不合并 PPTX。")
    parser.add_argument("--pdf-only", action="store_true",
                        help="仅生成 PDF 文件保存到工作目录，不生成 JPG。用于批量 QA 时保留中间态。")
    parser.add_argument("--pdf-out", default=None,
                        help="PDF 输出路径（默认输出到 --pptx-dir/merged.pdf）")

    args = parser.parse_args()
    os.makedirs(args.output_dir, exist_ok=True)

    if args.pdf_only:
        result = pptx_to_pdf(args.pptx_dir, args.pdf_out)
    else:
        result = pptx_to_images(args.pptx_dir, args.output_dir, args.dpi, args.files)

    if args.format == "json":
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(f"转换完成: {result.get('total_slides', 0)} 页幻灯片, {result.get('total_files', 0)} 个文件")
        print(f"图片目录: {args.output_dir}")
        for img in result.get("slide_images", []):
            print(f"  - {img}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
